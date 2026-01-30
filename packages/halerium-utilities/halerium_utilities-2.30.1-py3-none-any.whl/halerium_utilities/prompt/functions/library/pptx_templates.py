from pydantic.v1 import BaseModel, Field


def _find_substrings(text):
    """
    Find placeholder substrings in text
    """
    import re

    # Define the regex pattern to find substrings between {{ and }}
    pattern = r'\{\{.*?\}\}'

    # Use re.findall() to find all matches in the text
    matches = re.findall(pattern, text)

    return matches


def _collect_placeholders_from_item(text_frames, placeholders):
    """
    Get all placeholders from text frames
    """
    for text_frame in text_frames:
        for paragraph in text_frame.paragraphs:
            full_text = ''.join(run.text for run in paragraph.runs)
            matches = _find_substrings(full_text)
            placeholders.extend(matches)
    return placeholders


class ReturnPlaceholdersArgs(BaseModel):
    template_path: str = Field(description="The path to the pptx template.")


async def get_pptx_placeholders(data: ReturnPlaceholdersArgs):
    """
    Return list of placeholders of the pattern '{{key}}' in the pptx template.
    """

    import json
    from pptx import Presentation
    import asyncio

    template_path = data.get("template_path", "")

    prs = await asyncio.to_thread(Presentation, template_path)
    placeholders = []

    # look for placeholders in slides
    for i, slide in enumerate(prs.slides):
        slide_placeholders = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                slide_placeholders = _collect_placeholders_from_item([shape.text_frame], slide_placeholders)
        placeholders += slide_placeholders

    return str(placeholders)


async def _replace_text_preserving_formatting(text_frame, search_text, replacement_text):
    # Join all runs to handle placeholders spanning multiple runs
    for paragraph in text_frame.paragraphs:
        full_text = ''.join(run.text for run in paragraph.runs)
        if search_text in full_text:
            # Replace the text in the full string
            new_text = full_text.replace(search_text, replacement_text)
            # Clear the existing runs
            for run in paragraph.runs:
                run.text = ''
            # Add the new text back into the runs
            paragraph.runs[0].text = new_text


class FillTemplateArgs(BaseModel):
    template_path: str = Field(description="The path to the PPTX template.")
    output_path: str = Field(description="The path of the output file.")
    context: str = Field(
        description="The JSON string of the dictionary containing the placeholders as keys and their corresponding values to be replaced in the template. Multiple placeholders can be filled in a single call.")


async def fill_pptx_template(data: FillTemplateArgs):
    """
    This function creates a copy of the template with the placeholders replaced by the context.
    The context parameter should be a JSON string containing all the placeholders as keys and their corresponding values to be replaced in the template.

    Note:
    - For in-place updates, ensure that the output_path is the same as the template_path to maintain previously filled placeholders.
    - Example for context: '{"{{Title}}": "Hello", "{{Body}}": "world"}'.
    """

    import json
    from pptx import Presentation
    import asyncio

    template_path = data.get("template_path", "")
    output_path = data.get("output_path", "")
    context = data.get("context", "{}")
    context = json.loads(context)

    # Load the template presentation
    prs = await asyncio.to_thread(Presentation, template_path)

    # Replace placeholders in the presentation slides while preserving formatting
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for key, value in context.items():
                    await _replace_text_preserving_formatting(shape.text_frame, key, value)

    # Save the modified presentation
    await asyncio.to_thread(prs.save, output_path)
    return "Done"
