from typing import Union, IO

from pptx import Presentation

from .slide import DeckSlide

class Deck:
    """Manage a PowerPoint presentation."""

    def __init__(self, powerpoint: Union[str, IO[bytes]]):
        """
        Create a deck from a PowerPoint file path or binary stream.

        Args:
            powerpoint: Path to a PowerPoint file or a binary stream.
        """
        self.deck = Presentation(powerpoint)
        self.slides = self.deck.slides
        self._init_all_slides()

    def _init_all_slides(self) -> None:
        """Wrap all slides with DeckSlide."""
        self.slides = [DeckSlide(slide) for slide in self.slides]

    def get_slide(self, index: int) -> DeckSlide:
        """Return the slide at the given zero-based index."""

        try :
            slide = self.slides[index]
            return slide
            
        except IndexError:
            raise IndexError(f"Slide at index {index} not found")

    def save(self, path: str) -> None:
        """Save the deck to the given file path."""
        self.deck.save(path)
