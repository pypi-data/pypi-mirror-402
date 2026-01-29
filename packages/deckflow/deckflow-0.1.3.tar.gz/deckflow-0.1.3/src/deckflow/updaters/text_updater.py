from typing import Any

from pptx.dml.color import RGBColor

class TextUpdater:
    """Updater for text elements."""

    def __init__(self, color_analyzer: Any, formatter: Any):
        self.color_analyzer = color_analyzer
        self.formatter = formatter

    def update_text_frame(self, text_frame: Any, new_content: str, color_by_value: bool = False) -> None:
        # Find a template font from existing runs
        template_font = None
        for para in text_frame.paragraphs:
            for run in para.runs:
                if run.text and run.text.strip():
                    template_font = run.font
                    break
            if template_font:
                break

        lines = new_content.split("\n")

        # Remove all paragraphs except the first
        while len(text_frame.paragraphs) > 1:
            p = text_frame.paragraphs[-1]
            p._element.getparent().remove(p._element)

        # Empty the first paragraph
        text_frame.paragraphs[0].clear()

        # Add new lines with formatting
        for i, line in enumerate(lines):
            para = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
            run = para.add_run()
            run.text = line

            if template_font:
                self.formatter(template_font, run.font)

            if color_by_value:
                color_rgb = self.color_analyzer.get_color_for_value(line)
                if color_rgb:
                    run.font.color.rgb = RGBColor(*color_rgb)