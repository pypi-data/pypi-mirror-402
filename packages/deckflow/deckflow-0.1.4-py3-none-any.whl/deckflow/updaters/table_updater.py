from typing import Any, List
from pptx.util import Pt
from pptx.dml.color import RGBColor

class TableUpdater:
    """Updater for table elements."""

    def __init__(self, table: Any, color_analyzer: Any, formatter: Any):
        """
        table: DeckTable instance
        """
        self.table = table
        self.rows = table.rows
        self.cols = table.cols
        self.color_analyzer = color_analyzer
        self.formatter = formatter

    def update_cell(self, row: int, col: int, value: Any) -> bool:
        if 0 <= row < self.rows and 0 <= col < self.cols:
            self.table.data[row][col] = "" if value is None else str(value)
            print(f"Cell ({row},{col}) -> '{self.table.data[row][col]}'")
            return True
        print(f"Invalid cell indices ({row},{col}) for table {self.rows}x{self.cols}")
        return False

    def update_row(self, row_index: int, values: List[Any]) -> bool:
        if not (0 <= row_index < self.rows):
            print(f"Invalid row index {row_index}")
            return False
        for c in range(min(self.cols, len(values))):
            self.table.data[row_index][c] = "" if values[c] is None else str(values[c])
        print(f"Row {row_index} updated")
        return True

    def update_column(self, col_index: int, values: List[Any]) -> bool:
        if not (0 <= col_index < self.cols):
            print(f"Invalid column index {col_index}")
            return False
        for r in range(min(self.rows, len(values))):
            self.table.data[r][col_index] = "" if values[r] is None else str(values[r])
        print(f"Column {col_index} updated")
        return True
    
    def save_changes(self, color_by_value: bool = False):
        try:
            # Check if there are changes to save
            if self.table.data == self.table.original_data:
                print(f"No changes to save for table {self.table.name}")
                return True

            # Apply changes to the actual PowerPoint table object
            for row_idx, row in enumerate(self.table.table.rows):
                for col_idx, cell in enumerate(row.cells):
                    new_text = self.table.data[row_idx][col_idx]
                    if cell.text != new_text:
                        try:
                            if cell.text_frame and cell.text_frame.paragraphs:
                                para = cell.text_frame.paragraphs[0]
                                
                                # Extract font properties from the existing run BEFORE clearing
                                old_font = None
                                if para.runs:
                                    old_font = para.runs[0].font
                                
                                # Clear and add new run
                                para.clear()
                                run = para.add_run()
                                run.text = new_text

                                # Apply formatting from old font
                                if old_font:
                                    self.formatter(old_font, run.font)
                                
                                # Apply color by value if needed
                                if color_by_value:
                                    color_rgb = self.color_analyzer.get_color_for_value(new_text)
                                    if color_rgb:
                                        run.font.color.rgb = RGBColor(*color_rgb)

                            else:
                                cell.text = new_text
                        except Exception as e:
                            print(f"Warning: Could not update cell ({row_idx},{col_idx}): {e}")

            # Update the original state
            self.table.original_data = [row[:] for row in self.table.data]
            print(f"Table {self.table.name} changes applied")
            return True

        except Exception as e:
            print(f"Error saving table changes: {e}")
            return False