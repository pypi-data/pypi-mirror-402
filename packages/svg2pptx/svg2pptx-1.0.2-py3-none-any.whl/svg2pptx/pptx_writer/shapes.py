"""PowerPoint shape creation utilities."""

from typing import Optional

from pptx.shapes.base import BaseShape
from pptx.shapes.shapetree import SlideShapes, GroupShapes
from pptx.util import Emu
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

from svg2pptx.parser.styles import Style
from svg2pptx.parser.shapes import (
    ParsedShape,
    RectShape,
    CircleShape,
    EllipseShape,
    LineShape,
    PolygonShape,
    PolylineShape,
)
from svg2pptx.geometry.units import px_to_emu
from svg2pptx.geometry.transforms import Transform


def create_shape(
    shapes: SlideShapes,
    parsed_shape: ParsedShape,
    offset_x: int = 0,
    offset_y: int = 0,
    scale: float = 1.0,
) -> Optional[BaseShape]:
    """
    Create a PowerPoint shape from a parsed SVG shape.

    Args:
        shapes: SlideShapes or GroupShapes collection to add shape to.
        parsed_shape: Parsed shape data.
        offset_x: X offset in EMU.
        offset_y: Y offset in EMU.
        scale: Scale factor.

    Returns:
        Created shape or None.
    """
    if isinstance(parsed_shape, RectShape):
        return create_rectangle(shapes, parsed_shape, offset_x, offset_y, scale)
    elif isinstance(parsed_shape, (CircleShape, EllipseShape)):
        return create_oval(shapes, parsed_shape, offset_x, offset_y, scale)
    elif isinstance(parsed_shape, LineShape):
        return create_line(shapes, parsed_shape, offset_x, offset_y, scale)
    elif isinstance(parsed_shape, (PolygonShape, PolylineShape)):
        from svg2pptx.pptx_writer.freeform import create_freeform

        return create_freeform(shapes, parsed_shape, offset_x, offset_y, scale)
    return None


def create_rectangle(
    shapes: SlideShapes,
    rect: RectShape,
    offset_x: int = 0,
    offset_y: int = 0,
    scale: float = 1.0,
) -> BaseShape:
    """Create a PowerPoint rectangle shape."""
    # Apply transform to get actual position
    x, y = rect.transform.apply(rect.x, rect.y)
    
    # Convert to EMU with scale
    left = offset_x + px_to_emu(x * scale)
    top = offset_y + px_to_emu(y * scale)
    width = px_to_emu(rect.width * scale)
    height = px_to_emu(rect.height * scale)

    # Choose shape type based on corner radius
    if rect.rx > 0 or rect.ry > 0:
        shape = shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
        )
        # Note: python-pptx doesn't easily support setting corner radius
    else:
        shape = shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)

    apply_style(shape, rect.style)
    return shape


def create_oval(
    shapes: SlideShapes,
    oval: CircleShape | EllipseShape,
    offset_x: int = 0,
    offset_y: int = 0,
    scale: float = 1.0,
) -> BaseShape:
    """Create a PowerPoint oval (ellipse/circle) shape."""
    if isinstance(oval, CircleShape):
        cx, cy = oval.transform.apply(oval.cx, oval.cy)
        rx = ry = oval.r
    else:
        cx, cy = oval.transform.apply(oval.cx, oval.cy)
        rx, ry = oval.rx, oval.ry

    # Convert center + radius to left, top, width, height
    left = offset_x + px_to_emu((cx - rx) * scale)
    top = offset_y + px_to_emu((cy - ry) * scale)
    width = px_to_emu(2 * rx * scale)
    height = px_to_emu(2 * ry * scale)

    shape = shapes.add_shape(MSO_SHAPE.OVAL, left, top, width, height)
    apply_style(shape, oval.style)
    return shape


def create_line(
    shapes: SlideShapes,
    line: LineShape,
    offset_x: int = 0,
    offset_y: int = 0,
    scale: float = 1.0,
) -> BaseShape:
    """Create a PowerPoint line connector."""
    from pptx.enum.shapes import MSO_CONNECTOR

    # Apply transform
    x1, y1 = line.transform.apply(line.x1, line.y1)
    x2, y2 = line.transform.apply(line.x2, line.y2)

    # Convert to EMU
    start_x = offset_x + px_to_emu(x1 * scale)
    start_y = offset_y + px_to_emu(y1 * scale)
    end_x = offset_x + px_to_emu(x2 * scale)
    end_y = offset_y + px_to_emu(y2 * scale)

    connector = shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, start_x, start_y, end_x, end_y
    )

    # Disable shadow on line
    try:
        connector.shadow.inherit = False
        if hasattr(connector.shadow, 'visible'):
            connector.shadow.visible = False
    except (AttributeError, NotImplementedError):
        pass

    # Apply stroke style to line
    if line.style.stroke != "none":
        try:
            color = parse_hex_color(line.style.stroke)
            connector.line.color.rgb = color
        except ValueError:
            pass

    connector.line.width = Emu(px_to_emu(line.style.stroke_width))

    return connector


def apply_style(shape: BaseShape, style: Style, disable_shadow: bool = True) -> None:
    """
    Apply SVG style to a PowerPoint shape.

    Args:
        shape: PowerPoint shape to style.
        style: Parsed SVG style.
        disable_shadow: Whether to disable shadow on the shape. Defaults to True.
    """
    # Disable shadow if requested
    if disable_shadow:
        try:
            shape.shadow.inherit = False
            # Setting shadow to no shadow by making it transparent
            if hasattr(shape.shadow, 'visible'):
                shape.shadow.visible = False
        except (AttributeError, NotImplementedError):
            # Some shapes may not support shadow property
            pass

    # Apply fill
    fill = shape.fill
    if style.fill == "none":
        fill.background()  # No fill
    else:
        try:
            color = parse_hex_color(style.fill)
            fill.solid()
            fill.fore_color.rgb = color
            
            # Apply fill opacity
            if style.effective_fill_opacity < 1.0:
                # python-pptx doesn't directly support fill opacity
                # We'd need to modify the XML directly for this
                pass
        except ValueError:
            fill.background()

    # Apply stroke
    line = shape.line
    if style.stroke == "none":
        line.fill.background()  # No stroke
    else:
        try:
            color = parse_hex_color(style.stroke)
            line.color.rgb = color
            line.width = Emu(px_to_emu(style.stroke_width))
        except ValueError:
            line.fill.background()


def parse_hex_color(hex_color: str) -> RGBColor:
    """
    Parse a hex color string to RGBColor.

    Args:
        hex_color: Color in format "#RRGGBB" or "#RGB".

    Returns:
        RGBColor object.

    Raises:
        ValueError: If color format is invalid.
    """
    if not hex_color or hex_color == "none":
        raise ValueError("Invalid color: none")

    color = hex_color.strip().lstrip("#")

    if len(color) == 3:
        color = "".join(c * 2 for c in color)

    if len(color) != 6:
        raise ValueError(f"Invalid hex color: {hex_color}")

    try:
        r = int(color[0:2], 16)
        g = int(color[2:4], 16)
        b = int(color[4:6], 16)
        return RGBColor(r, g, b)
    except ValueError:
        raise ValueError(f"Invalid hex color: {hex_color}")
