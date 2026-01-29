"""PowerPoint freeform shape creation for polygons, polylines, and paths."""

from typing import Optional, Union

from pptx.shapes.base import BaseShape
from pptx.shapes.shapetree import SlideShapes, GroupShapes
from pptx.util import Emu

from svg2pptx.parser.shapes import PolygonShape, PolylineShape
from svg2pptx.parser.paths import PathShape
from svg2pptx.parser.styles import Style
from svg2pptx.geometry.units import px_to_emu
from svg2pptx.pptx_writer.shapes import apply_style


def create_freeform(
    shapes: SlideShapes,
    parsed_shape: Union[PolygonShape, PolylineShape, PathShape],
    offset_x: int = 0,
    offset_y: int = 0,
    scale: float = 1.0,
) -> Optional[BaseShape]:
    """
    Create a PowerPoint freeform shape from polygon, polyline, or path.

    Args:
        shapes: SlideShapes collection to add shape to.
        parsed_shape: Parsed polygon, polyline, or path shape.
        offset_x: X offset in EMU.
        offset_y: Y offset in EMU.
        scale: Scale factor.

    Returns:
        Created freeform shape or None.
    """
    if isinstance(parsed_shape, PathShape):
        return create_freeform_from_path(
            shapes, parsed_shape, offset_x, offset_y, scale
        )
    else:
        return create_freeform_from_points(
            shapes, parsed_shape, offset_x, offset_y, scale
        )


def create_freeform_from_points(
    shapes: SlideShapes,
    parsed_shape: Union[PolygonShape, PolylineShape],
    offset_x: int = 0,
    offset_y: int = 0,
    scale: float = 1.0,
) -> Optional[BaseShape]:
    """
    Create a freeform shape from polygon or polyline points.

    Args:
        shapes: SlideShapes collection.
        parsed_shape: Parsed polygon or polyline.
        offset_x: X offset in EMU.
        offset_y: Y offset in EMU.
        scale: Scale factor.

    Returns:
        Created shape or None.
    """
    points = parsed_shape.points
    if len(points) < 2:
        return None

    # Apply transform to all points
    transformed_points = parsed_shape.transform.apply_to_points(points)

    # Convert to EMU and apply scale
    emu_points = [
        (px_to_emu(x * scale), px_to_emu(y * scale))
        for x, y in transformed_points
    ]

    # Determine if closed
    is_closed = isinstance(parsed_shape, PolygonShape)

    # Create freeform using FreeformBuilder
    first_x, first_y = emu_points[0]
    builder = shapes.build_freeform(first_x, first_y)
    builder.add_line_segments(emu_points[1:], close=is_closed)

    # Convert to shape
    shape = builder.convert_to_shape(offset_x, offset_y)

    # Apply styling
    apply_style(shape, parsed_shape.style)

    return shape


def create_freeform_from_path(
    shapes: SlideShapes,
    path: PathShape,
    offset_x: int = 0,
    offset_y: int = 0,
    scale: float = 1.0,
) -> Optional[BaseShape]:
    """
    Create a freeform shape from an SVG path.

    Paths with multiple subpaths are combined into a single shape.

    Args:
        shapes: SlideShapes collection.
        path: Parsed path shape.
        offset_x: X offset in EMU.
        offset_y: Y offset in EMU.
        scale: Scale factor.

    Returns:
        Created shape or None.
    """
    if not path.subpaths:
        return None

    # Get the first subpath to start the builder
    first_subpath_points, first_closed = path.subpaths[0]
    if len(first_subpath_points) < 2:
        return None

    # Apply transform to first point
    transformed_first = path.transform.apply_to_points(first_subpath_points)
    first_x = px_to_emu(transformed_first[0][0] * scale)
    first_y = px_to_emu(transformed_first[0][1] * scale)

    # Start the freeform builder
    builder = shapes.build_freeform(first_x, first_y)

    # Add segments for first subpath
    emu_points = [
        (px_to_emu(x * scale), px_to_emu(y * scale))
        for x, y in transformed_first[1:]
    ]
    builder.add_line_segments(emu_points, close=first_closed)

    # Add additional subpaths using move_to
    for subpath_points, is_closed in path.subpaths[1:]:
        if len(subpath_points) < 2:
            continue

        transformed = path.transform.apply_to_points(subpath_points)
        
        # Move to start of new subpath
        start_x = px_to_emu(transformed[0][0] * scale)
        start_y = px_to_emu(transformed[0][1] * scale)
        builder.move_to(start_x, start_y)

        # Add line segments
        emu_points = [
            (px_to_emu(x * scale), px_to_emu(y * scale))
            for x, y in transformed[1:]
        ]
        builder.add_line_segments(emu_points, close=is_closed)

    # Convert to shape
    shape = builder.convert_to_shape(offset_x, offset_y)

    # Apply styling
    apply_style(shape, path.style)

    return shape
