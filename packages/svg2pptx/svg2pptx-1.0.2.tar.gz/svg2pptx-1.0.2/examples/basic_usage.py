"""Example: Basic usage of svg2pptx."""

from pathlib import Path
from svg2pptx import svg_to_pptx, SVGConverter, Config

# Directory containing this script
SCRIPT_DIR = Path(__file__).parent
PROJECT_DIR = SCRIPT_DIR.parent


def example_simple_conversion():
    """Simple one-line conversion."""
    print("Example 1: Simple conversion")
    
    svg_path = PROJECT_DIR / "tests" / "fixtures" / "basic_shapes.svg"
    output_path = PROJECT_DIR / "examples" / "output_simple.pptx"
    
    svg_to_pptx(str(svg_path), str(output_path))
    print(f"  Created: {output_path}")


def example_with_config():
    """Conversion with custom configuration."""
    print("\nExample 2: Conversion with config")
    
    svg_path = PROJECT_DIR / "tests" / "fixtures" / "path_icon.svg"
    output_path = PROJECT_DIR / "examples" / "output_scaled.pptx"
    
    config = Config(
        scale=1.5,                  # Scale up by 50%
        curve_tolerance=0.5,        # Higher quality curves
        preserve_groups=True,       # Keep group structure
    )
    
    svg_to_pptx(str(svg_path), str(output_path), config=config)
    print(f"  Created: {output_path}")


def example_add_to_existing():
    """Add SVG to an existing presentation."""
    print("\nExample 3: Add to existing presentation")
    
    from pptx import Presentation
    from pptx.util import Inches
    
    # Create a new presentation
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Add a slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Add SVG content
    converter = SVGConverter()
    svg_path = PROJECT_DIR / "tests" / "fixtures" / "path_icon.svg"
    converter.add_to_slide(str(svg_path), slide)
    
    output_path = PROJECT_DIR / "examples" / "output_existing.pptx"
    prs.save(str(output_path))
    print(f"  Created: {output_path}")


def example_from_string():
    """Convert SVG from a string."""
    print("\nExample 4: From SVG string")
    
    svg_content = '''<svg xmlns="http://www.w3.org/2000/svg" width="200" height="100">
        <rect x="10" y="10" width="80" height="80" fill="#3498db" rx="10"/>
        <circle cx="150" cy="50" r="40" fill="#e74c3c"/>
    </svg>'''
    
    converter = SVGConverter()
    prs = converter.convert_string(svg_content)
    
    output_path = PROJECT_DIR / "examples" / "output_from_string.pptx"
    prs.save(str(output_path))
    print(f"  Created: {output_path}")


if __name__ == "__main__":
    print("SVG to PowerPoint Converter - Examples\n")
    print("=" * 50)
    
    # Ensure output directory exists
    output_dir = PROJECT_DIR / "examples"
    output_dir.mkdir(exist_ok=True)
    
    try:
        example_simple_conversion()
        example_with_config()
        example_add_to_existing()
        example_from_string()
        
        print("\n" + "=" * 50)
        print("All examples completed successfully!")
        print(f"Output files are in: {output_dir}")
    except Exception as e:
        print(f"\nError: {e}")
        raise
