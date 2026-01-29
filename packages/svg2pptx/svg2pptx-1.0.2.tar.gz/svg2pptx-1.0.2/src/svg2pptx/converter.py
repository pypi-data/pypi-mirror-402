"""Main SVG to PowerPoint converter."""

from pathlib import Path
from typing import Optional, Union

from pptx import Presentation
from pptx.slide import Slide
from pptx.util import Emu

from svg2pptx.config import Config
from svg2pptx.parser.svg_parser import SVGParser, SVGDocument, GroupElement, TextElement
from svg2pptx.parser.shapes import ParsedShape
from svg2pptx.parser.paths import PathShape
from svg2pptx.geometry.units import px_to_emu
from svg2pptx.pptx_writer.shapes import create_shape
from svg2pptx.pptx_writer.freeform import create_freeform
from svg2pptx.pptx_writer.groups import create_group, add_element_to_shapes
from svg2pptx.pptx_writer.text import create_text


class SVGConverter:
    """
    Converts SVG files to PowerPoint presentations.

    Supports basic shapes, paths, groups, and text elements.
    """

    def __init__(self, config: Optional[Config] = None):
        """
        Initialize the SVG converter.

        Args:
            config: Optional configuration settings.
        """
        self.config = config or Config()
        self.parser = SVGParser(curve_tolerance=self.config.curve_tolerance)

    def convert_file(
        self,
        svg_path: Union[str, Path],
        pptx_path: Union[str, Path],
    ) -> None:
        """
        Convert an SVG file to a PowerPoint presentation.

        Args:
            svg_path: Path to the input SVG file.
            pptx_path: Path for the output PowerPoint file.
        """
        # Parse the SVG
        svg_doc = self.parser.parse_file(svg_path)

        # Create presentation
        prs = self._create_presentation(svg_doc)

        # Save
        prs.save(str(pptx_path))

    def convert_string(self, svg_content: str) -> Presentation:
        """
        Convert an SVG string to a PowerPoint Presentation object.

        Args:
            svg_content: SVG content as a string.

        Returns:
            PowerPoint Presentation object.
        """
        svg_doc = self.parser.parse_string(svg_content)
        return self._create_presentation(svg_doc)

    def add_to_slide(
        self,
        svg_path: Union[str, Path],
        slide: Slide,
        x: Optional[int] = None,
        y: Optional[int] = None,
        scale: Optional[float] = None,
    ) -> None:
        """
        Add SVG shapes to an existing slide.

        Args:
            svg_path: Path to the SVG file.
            slide: PowerPoint slide to add shapes to.
            x: X position in EMU (default: use config offset_x).
            y: Y position in EMU (default: use config offset_y).
            scale: Scale factor (default: use config scale).
        """
        svg_doc = self.parser.parse_file(svg_path)
        self._add_svg_to_slide(
            svg_doc,
            slide,
            x if x is not None else self.config.offset_x,
            y if y is not None else self.config.offset_y,
            scale if scale is not None else self.config.scale,
        )

    def add_string_to_slide(
        self,
        svg_content: str,
        slide: Slide,
        x: Optional[int] = None,
        y: Optional[int] = None,
        scale: Optional[float] = None,
    ) -> None:
        """
        Add SVG shapes from a string to an existing slide.

        Args:
            svg_content: SVG content as a string.
            slide: PowerPoint slide to add shapes to.
            x: X position in EMU.
            y: Y position in EMU.
            scale: Scale factor.
        """
        svg_doc = self.parser.parse_string(svg_content)
        self._add_svg_to_slide(
            svg_doc,
            slide,
            x if x is not None else self.config.offset_x,
            y if y is not None else self.config.offset_y,
            scale if scale is not None else self.config.scale,
        )

    def _create_presentation(self, svg_doc: SVGDocument) -> Presentation:
        """Create a new presentation and add the SVG content."""
        prs = Presentation()

        # Set slide dimensions
        prs.slide_width = Emu(self.config.slide_width)
        prs.slide_height = Emu(self.config.slide_height)

        # Use blank layout (index 6)
        try:
            blank_layout = prs.slide_layouts[6]
        except IndexError:
            blank_layout = prs.slide_layouts[0]

        slide = prs.slides.add_slide(blank_layout)

        # Calculate scale to fit SVG into slide
        scale = self._calculate_scale(svg_doc)

        # Add SVG content
        self._add_svg_to_slide(
            svg_doc,
            slide,
            self.config.offset_x,
            self.config.offset_y,
            scale * self.config.scale,
        )

        return prs

    def _calculate_scale(self, svg_doc: SVGDocument) -> float:
        """Calculate scale factor to fit SVG into slide."""
        if svg_doc.width == 0 or svg_doc.height == 0:
            return 1.0

        # Calculate scale to fit
        svg_width_emu = px_to_emu(svg_doc.width)
        svg_height_emu = px_to_emu(svg_doc.height)

        scale_x = self.config.slide_width / svg_width_emu
        scale_y = self.config.slide_height / svg_height_emu

        # Use smaller scale to fit completely
        return min(scale_x, scale_y, 1.0)

    def _add_svg_to_slide(
        self,
        svg_doc: SVGDocument,
        slide: Slide,
        offset_x: int,
        offset_y: int,
        scale: float,
    ) -> None:
        """Add SVG elements to a slide."""
        # Apply viewBox transform if present
        vb_scale_x = svg_doc.scale_x * scale
        vb_scale_y = svg_doc.scale_y * scale
        vb_offset_x = offset_x + px_to_emu(svg_doc.offset_x * scale)
        vb_offset_y = offset_y + px_to_emu(svg_doc.offset_y * scale)

        # Use average scale for uniform scaling
        effective_scale = (vb_scale_x + vb_scale_y) / 2

        for element in svg_doc.elements:
            add_element_to_shapes(
                slide.shapes,
                element,
                vb_offset_x,
                vb_offset_y,
                effective_scale,
                flatten=self.config.flatten_groups,
                config=self.config,
            )
