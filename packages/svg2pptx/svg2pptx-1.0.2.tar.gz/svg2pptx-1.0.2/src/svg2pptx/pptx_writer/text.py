"""PowerPoint text box creation from SVG text elements."""

from typing import Optional

from pptx.shapes.base import BaseShape
from pptx.shapes.shapetree import SlideShapes
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from svg2pptx.parser.svg_parser import TextElement
from svg2pptx.geometry.units import px_to_emu
from svg2pptx.pptx_writer.shapes import parse_hex_color


def create_text(
    shapes: SlideShapes,
    text_element: TextElement,
    offset_x: int = 0,
    offset_y: int = 0,
    scale: float = 1.0,
) -> Optional[BaseShape]:
    """
    Create a PowerPoint text box from an SVG text element.

    Args:
        shapes: SlideShapes collection to add text box to.
        text_element: Parsed SVG text element.
        offset_x: X offset in EMU.
        offset_y: Y offset in EMU.
        scale: Scale factor.

    Returns:
        Created text box shape or None.
    """
    if not text_element.text:
        return None

    # Apply transform to position
    x, y = text_element.transform.apply(text_element.x, text_element.y)

    # Get font size in pixels
    font_size_px = text_element.style.font_size * scale
    
    # Estimate text box dimensions
    # Use a more generous width estimate to avoid text wrapping
    estimated_width = len(text_element.text) * font_size_px * 0.7 + font_size_px
    estimated_height = font_size_px * 1.4

    # Convert position to EMU
    # In SVG, text y-coordinate is the baseline position
    # We need to adjust based on the text-anchor for horizontal positioning
    text_anchor = text_element.style.text_anchor
    
    # Calculate left position based on text-anchor
    if text_anchor == "middle":
        # Text is centered at x
        left = offset_x + px_to_emu(x * scale) - px_to_emu(estimated_width / 2)
    elif text_anchor == "end":
        # Text ends at x
        left = offset_x + px_to_emu(x * scale) - px_to_emu(estimated_width)
    else:
        # Default: text starts at x (text-anchor="start")
        left = offset_x + px_to_emu(x * scale)

    # For y position: SVG y is the baseline, so we move up by approximately
    # the font ascent (roughly 80% of font size for most fonts)
    baseline_offset = font_size_px * 0.85
    top = offset_y + px_to_emu(y * scale) - px_to_emu(baseline_offset)
    
    width = px_to_emu(estimated_width)
    height = px_to_emu(estimated_height)

    # Create text box
    text_box = shapes.add_textbox(left, top, width, height)
    text_frame = text_box.text_frame
    text_frame.word_wrap = False
    
    # Remove margins/padding for more accurate positioning
    text_frame.margin_left = 0
    text_frame.margin_right = 0
    text_frame.margin_top = 0
    text_frame.margin_bottom = 0

    # Set text content
    paragraph = text_frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = text_element.text

    # Apply text styling
    font = run.font
    font.name = text_element.style.font_family
    font.size = Pt(text_element.style.font_size * scale)

    # Font weight
    if text_element.style.font_weight in ("bold", "700", "800", "900"):
        font.bold = True

    # Text color (use fill color for text)
    if text_element.style.fill != "none":
        try:
            color = parse_hex_color(text_element.style.fill)
            font.color.rgb = color
        except ValueError:
            pass

    # Text anchor (horizontal alignment)
    anchor_map = {
        "start": PP_ALIGN.LEFT,
        "middle": PP_ALIGN.CENTER,
        "end": PP_ALIGN.RIGHT,
    }
    paragraph.alignment = anchor_map.get(text_anchor, PP_ALIGN.LEFT)

    # Disable shadow on text box
    try:
        text_box.shadow.inherit = False
        if hasattr(text_box.shadow, 'visible'):
            text_box.shadow.visible = False
    except (AttributeError, NotImplementedError):
        pass

    return text_box

