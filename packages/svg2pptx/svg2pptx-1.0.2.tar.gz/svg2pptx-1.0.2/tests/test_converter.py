"""Tests for the main SVG to PPTX converter."""

import pytest
from pathlib import Path
import tempfile
import os

from pptx import Presentation

from svg2pptx import svg_to_pptx, SVGConverter, Config


FIXTURES_DIR = Path(__file__).parent / "fixtures"


class TestSVGConverter:
    """Tests for SVGConverter class."""

    def test_convert_simple_svg_string(self):
        """Test converting a simple SVG string."""
        svg = '''<svg xmlns="http://www.w3.org/2000/svg" width="100" height="100">
            <rect x="10" y="10" width="80" height="80" fill="red"/>
        </svg>'''
        
        converter = SVGConverter()
        prs = converter.convert_string(svg)
        
        assert prs is not None
        assert len(prs.slides) == 1
        # Should have at least one shape
        assert len(prs.slides[0].shapes) >= 1

    def test_convert_circle(self):
        """Test converting a circle."""
        svg = '''<svg xmlns="http://www.w3.org/2000/svg" width="100" height="100">
            <circle cx="50" cy="50" r="40" fill="blue"/>
        </svg>'''
        
        converter = SVGConverter()
        prs = converter.convert_string(svg)
        
        assert len(prs.slides[0].shapes) >= 1

    def test_convert_with_config(self):
        """Test conversion with custom config."""
        svg = '''<svg xmlns="http://www.w3.org/2000/svg" width="100" height="100">
            <rect x="0" y="0" width="50" height="50"/>
        </svg>'''
        
        config = Config(scale=2.0)
        converter = SVGConverter(config=config)
        prs = converter.convert_string(svg)
        
        assert prs is not None

    def test_convert_multiple_shapes(self):
        """Test converting multiple shapes."""
        svg = '''<svg xmlns="http://www.w3.org/2000/svg" width="200" height="100">
            <rect x="10" y="10" width="50" height="50" fill="red"/>
            <circle cx="120" cy="35" r="25" fill="blue"/>
            <ellipse cx="170" cy="60" rx="20" ry="30" fill="green"/>
        </svg>'''
        
        converter = SVGConverter()
        prs = converter.convert_string(svg)
        
        # Should have 3 shapes
        assert len(prs.slides[0].shapes) >= 3


class TestConvertFile:
    """Tests for file-based conversion."""

    def test_convert_basic_shapes(self):
        """Test converting the basic shapes fixture."""
        svg_path = FIXTURES_DIR / "basic_shapes.svg"
        if not svg_path.exists():
            pytest.skip("Fixture file not found")
        
        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
            pptx_path = f.name
        
        try:
            svg_to_pptx(str(svg_path), pptx_path)
            
            # Verify the output file exists and is valid
            assert os.path.exists(pptx_path)
            prs = Presentation(pptx_path)
            assert len(prs.slides) == 1
            assert len(prs.slides[0].shapes) > 0
        finally:
            if os.path.exists(pptx_path):
                os.unlink(pptx_path)

    def test_convert_path_icon(self):
        """Test converting a path-based icon."""
        svg_path = FIXTURES_DIR / "path_icon.svg"
        if not svg_path.exists():
            pytest.skip("Fixture file not found")
        
        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
            pptx_path = f.name
        
        try:
            svg_to_pptx(str(svg_path), pptx_path)
            
            assert os.path.exists(pptx_path)
            prs = Presentation(pptx_path)
            assert len(prs.slides[0].shapes) > 0
        finally:
            if os.path.exists(pptx_path):
                os.unlink(pptx_path)

    def test_convert_grouped_svg(self):
        """Test converting SVG with groups."""
        svg_path = FIXTURES_DIR / "grouped.svg"
        if not svg_path.exists():
            pytest.skip("Fixture file not found")
        
        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
            pptx_path = f.name
        
        try:
            config = Config(preserve_groups=True)
            svg_to_pptx(str(svg_path), pptx_path, config=config)
            
            assert os.path.exists(pptx_path)
            prs = Presentation(pptx_path)
            assert len(prs.slides[0].shapes) > 0
        finally:
            if os.path.exists(pptx_path):
                os.unlink(pptx_path)


class TestAddToSlide:
    """Tests for adding SVG to existing slides."""

    def test_add_to_existing_slide(self):
        """Test adding SVG shapes to an existing presentation."""
        svg = '''<svg xmlns="http://www.w3.org/2000/svg" width="100" height="100">
            <rect x="0" y="0" width="50" height="50" fill="red"/>
        </svg>'''
        
        # Create a presentation with one slide
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Add SVG shapes
        converter = SVGConverter()
        converter.add_string_to_slide(svg, slide)
        
        # Should have at least one shape
        assert len(slide.shapes) >= 1

    def test_add_multiple_svgs(self):
        """Test adding multiple SVGs to the same slide."""
        svg1 = '''<svg xmlns="http://www.w3.org/2000/svg" width="50" height="50">
            <rect x="0" y="0" width="50" height="50" fill="red"/>
        </svg>'''
        svg2 = '''<svg xmlns="http://www.w3.org/2000/svg" width="50" height="50">
            <circle cx="25" cy="25" r="25" fill="blue"/>
        </svg>'''
        
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        converter = SVGConverter()
        converter.add_string_to_slide(svg1, slide, x=0, y=0)
        converter.add_string_to_slide(svg2, slide, x=500000, y=0)
        
        # Should have at least 2 shapes
        assert len(slide.shapes) >= 2


class TestConfig:
    """Tests for configuration options."""

    def test_default_config(self):
        config = Config()
        assert config.scale == 1.0
        assert config.preserve_groups is False
        assert config.flatten_groups is True
        assert config.disable_shadows is True

    def test_custom_config(self):
        config = Config(
            scale=2.0,
            curve_tolerance=0.5,
            preserve_groups=False,
        )
        assert config.scale == 2.0
        assert config.curve_tolerance == 0.5
        assert config.preserve_groups is False
