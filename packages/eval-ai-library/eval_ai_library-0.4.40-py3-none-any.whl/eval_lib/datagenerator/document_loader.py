# document_loader.py
from __future__ import annotations
from pathlib import Path
from typing import List
from langchain_core.documents import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_community.document_loaders import PyPDFLoader
from langchain_community.document_loaders import Docx2txtLoader
from langchain_community.document_loaders import TextLoader

import html2text
import markdown
from pptx import Presentation
from striprtf.striprtf import rtf_to_text
from pypdf import PdfReader
import fitz
import zipfile
from xml.etree import ElementTree as ET
import pandas as pd
import yaml
import pytesseract
from PIL import Image
import io as _io
import pytesseract
from PIL import Image
import io as _io
import docx  # python-docx
import mammoth
import io
import json
import zipfile

try:
    import textract
    HAS_TEXTRACT = True
except ImportError:
    HAS_TEXTRACT = False
    textract = None

# ---------------------------
# Helper functions
# ---------------------------


def _read_text(p: Path) -> str:
    return p.read_text(encoding="utf-8", errors="ignore")


def _read_bytes(p: Path) -> bytes:
    return p.read_bytes()


def _csv_tsv_to_text(p: Path) -> str:
    try:

        sep = "," if p.suffix.lower() == ".csv" else "\t"
        df = pd.read_csv(str(p), dtype=str, sep=sep,
                         encoding="utf-8", engine="python")
        df = df.fillna("")
        buf = io.StringIO()
        df.to_csv(buf, index=False)
        return buf.getvalue()
    except Exception:
        try:
            return _read_text(p)
        except Exception:
            return ""


def _xlsx_to_text(p: Path) -> str:
    try:
        df = pd.read_excel(str(p), dtype=str, engine="openpyxl")
        df = df.fillna("")
        buf = io.StringIO()
        df.to_csv(buf, index=False)
        return buf.getvalue()
    except Exception:
        return ""


def _pptx_to_text(p: Path) -> str:
    try:
        prs = Presentation(str(p))
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    texts.append(shape.text)
        return "\n".join(texts)
    except Exception:
        return ""


def _json_to_text(p: Path) -> str:
    try:
        data = json.loads(_read_text(p))

        def flatten(obj, prefix=""):
            lines = []
            if isinstance(obj, dict):
                for k, v in obj.items():
                    lines += flatten(v, f"{prefix}{k}.")
            elif isinstance(obj, list):
                for i, v in enumerate(obj):
                    lines += flatten(v, f"{prefix}{i}.")
            else:
                lines.append(f"{prefix[:-1]}: {obj}")
            return lines
        return "\n".join(flatten(data))
    except Exception:
        return _read_text(p)


def _yaml_to_text(p: Path) -> str:
    try:
        data = yaml.safe_load(_read_text(p))
        return json.dumps(data, ensure_ascii=False, indent=2)
    except Exception:
        return _read_text(p)


def _xml_to_text(p: Path) -> str:
    try:
        tree = ET.parse(str(p))
        root = tree.getroot()
        lines = []

        def walk(node, prefix=""):
            text = (node.text or "").strip()
            tag = node.tag
            if text:
                lines.append(f"{prefix}{tag}: {text}")
            for child in node:
                walk(child, prefix + tag + ".")
        walk(root)
        return "\n".join(lines)
    except Exception:
        return _read_text(p)


def _rtf_to_text(p: Path) -> str:
    try:

        return rtf_to_text(_read_text(p))
    except Exception:
        return ""


def _odt_to_text(p: Path) -> str:
    try:
        with zipfile.ZipFile(str(p)) as z:
            from xml.etree import ElementTree as ET
            with z.open("content.xml") as f:
                tree = ET.parse(f)
                root = tree.getroot()
                texts = []
                for elem in root.iter():
                    if elem.text and elem.text.strip():
                        texts.append(elem.text.strip())
                return "\n".join(texts)
    except Exception:
        return ""

# ---------------------------
# PDF: LangChain -> pypdf -> PyMuPDF -> OCR(PyMuPDF+pytesseract)
# ---------------------------


def _pdf_text_pypdf(p: Path) -> str:
    try:
        reader = PdfReader(str(p))
        texts = []
        for page in reader.pages:
            t = page.extract_text() or ""
            if t.strip():
                texts.append(t)
        return "\n".join(texts)
    except Exception:
        return ""


def _pdf_text_pymupdf(p: Path) -> str:
    try:
        text_parts = []
        with fitz.open(str(p)) as doc:
            for page in doc:
                t = page.get_text("text") or ""
                if t.strip():
                    text_parts.append(t)
        return "\n".join(text_parts)
    except Exception:
        return ""


def _pdf_ocr_via_pymupdf(p: Path) -> str:
    """Render pages via PyMuPDF and OCR pytesseract. Will work if pytesseract + tesseract are installed."""
    try:
        texts = []
        zoom = 2.0
        mat = fitz.Matrix(zoom, zoom)
        with fitz.open(str(p)) as doc:
            for page in doc:
                pix = page.get_pixmap(matrix=mat, alpha=False)
                img = Image.open(_io.BytesIO(pix.tobytes("png")))
                t = pytesseract.image_to_string(img) or ""
                if t.strip():
                    texts.append(t)
        return "\n".join(texts)
    except Exception:
        return ""

# ---------------------------
# Images (OCR)
# ---------------------------


def _ocr_image_bytes(img_bytes: bytes) -> str:
    try:
        img = Image.open(_io.BytesIO(img_bytes))
        return pytesseract.image_to_string(img) or ""
    except Exception:
        return ""

# ---------------------------
# Docx
# ---------------------------


def _docx_to_text_python_docx(p: Path) -> str:
    try:
        d = docx.Document(str(p))
        parts = []
        for para in d.paragraphs:
            if para.text:
                parts.append(para.text)
        # захватим текст из таблиц
        for tbl in d.tables:
            for row in tbl.rows:
                cells = [cell.text for cell in row.cells]
                if any(c.strip() for c in cells):
                    parts.append("\t".join(cells))
        return "\n".join(parts)
    except Exception:
        return ""


def _docx_to_text_mammoth(p: Path) -> str:
    try:
        with open(str(p), "rb") as f:
            result = mammoth.extract_raw_text(f)
            return (result.value or "").strip()
    except Exception:
        return ""


def _docx_to_text_zipxml(p: Path) -> str:
    """Без зависимостей: читаем word/document.xml и вытаскиваем все w:t."""
    try:

        texts = []
        with zipfile.ZipFile(str(p)) as z:
            if "word/document.xml" in z.namelist():
                with z.open("word/document.xml") as f:
                    root = ET.parse(f).getroot()
                    for el in root.iter():
                        tag = el.tag.rsplit("}", 1)[-1]
                        if tag == "t" and el.text and el.text.strip():
                            texts.append(el.text.strip())
            for name in z.namelist():
                if name.startswith("word/header") and name.endswith(".xml"):
                    with z.open(name) as f:
                        root = ET.parse(f).getroot()
                        for el in root.iter():
                            tag = el.tag.rsplit("}", 1)[-1]
                            if tag == "t" and el.text and el.text.strip():
                                texts.append(el.text.strip())
                if name.startswith("word/footer") and name.endswith(".xml"):
                    with z.open(name) as f:
                        root = ET.parse(f).getroot()
                        for el in root.iter():
                            tag = el.tag.rsplit("}", 1)[-1]
                            if tag == "t" and el.text and el.text.strip():
                                texts.append(el.text.strip())
        return "\n".join(texts)
    except Exception:
        return ""


def _doc_to_text_textract(p: Path) -> str:
    if not HAS_TEXTRACT:
        return ""
    try:
        return textract.process(str(p)).decode("utf-8", errors="ignore")
    except Exception:
        return ""

# ---------------------------
# General functions (extended)
# ---------------------------


def load_documents(file_paths: List[str]) -> List[Document]:
    documents: List[Document] = []

    for path in map(Path, file_paths):
        ext = path.suffix.lower()

        try:
            # ---- PDF ----
            if ext == ".pdf":
                used_langchain = False
                # 1) LangChain PyPDFLoader
                try:
                    docs = PyPDFLoader(str(path)).load()
                    if docs and any((d.page_content or "").strip() for d in docs):
                        documents += docs
                        used_langchain = True
                except Exception:
                    used_langchain = False

                if not used_langchain:
                    # 2) pypdf
                    text = _pdf_text_pypdf(path)
                    if not text.strip():
                        # 3) PyMuPDF
                        text = _pdf_text_pymupdf(path)
                    if not text.strip():
                        # 4) OCR через PyMuPDF
                        text = _pdf_ocr_via_pymupdf(path)

                    if text.strip():
                        documents.append(Document(page_content=text, metadata={
                                         "source": str(path), "filetype": "pdf"}))
                    else:
                        print(
                            f"⚠️  PDF has no extractable text (maybe scanned): {path.name}")

            # ---- DOCX ----
            elif ext == ".docx":
                # 1) Пытаемся стандартным Docx2txtLoader
                added = False
                try:
                    docs = Docx2txtLoader(str(path)).load()
                    # Docx2txt иногда возвращает Document с пустым page_content
                    docs = [d for d in docs if (d.page_content or "").strip()]
                    if docs:
                        documents += docs
                        added = True
                except Exception:
                    added = False

                if not added:
                    # 2) python-docx
                    text = _docx_to_text_python_docx(path)
                    if not text.strip():
                        # 3) mammoth
                        text = _docx_to_text_mammoth(path)
                    if not text.strip():
                        # 4) zip+xml fallback
                        text = _docx_to_text_zipxml(path)

                    if text.strip():
                        documents.append(Document(
                            page_content=text,
                            metadata={"source": str(path), "filetype": "docx"}
                        ))
                    else:
                        print(f"⚠️  DOCX produced no text: {path.name}")

            elif ext == ".doc":
                # старый формат
                text = _doc_to_text_textract(path)
                if text.strip():
                    documents.append(Document(
                        page_content=text,
                        metadata={"source": str(path), "filetype": "doc"}
                    ))
                else:
                    print(
                        f"⚠️  .DOC not extractable (install textract/antiword?): {path.name}")

            # ---- TXT ----
            elif ext == ".txt":
                documents += TextLoader(str(path), encoding="utf-8").load()

            # ---- HTML ----
            elif ext in (".html", ".htm"):
                html = _read_text(path)
                text = html2text.html2text(html)
                documents.append(Document(page_content=text, metadata={
                                 "source": str(path), "filetype": "html"}))

            # ---- Markdown ----
            elif ext == ".md":
                md = _read_text(path)
                html = markdown.markdown(md)
                text = html2text.html2text(html)
                documents.append(Document(page_content=text, metadata={
                                 "source": str(path), "filetype": "md"}))

            # ---- CSV / TSV ----
            elif ext in (".csv", ".tsv"):
                text = _csv_tsv_to_text(path)
                if text.strip():
                    documents.append(Document(page_content=text, metadata={
                                     "source": str(path), "filetype": ext.lstrip(".")}))
                else:
                    print(f"⚠️  Empty CSV/TSV: {path.name}")

            # ---- XLSX ----
            elif ext == ".xlsx":
                text = _xlsx_to_text(path)
                if text.strip():
                    documents.append(Document(page_content=text, metadata={
                                     "source": str(path), "filetype": "xlsx"}))
                else:
                    print(f"⚠️  Empty XLSX: {path.name}")

            # ---- PPTX ----
            elif ext == ".pptx":
                text = _pptx_to_text(path)
                if text.strip():
                    documents.append(Document(page_content=text, metadata={
                                     "source": str(path), "filetype": "pptx"}))
                else:
                    print(f"⚠️  Empty PPTX: {path.name}")

            # ---- JSON ----
            elif ext == ".json":
                text = _json_to_text(path)
                if text.strip():
                    documents.append(Document(page_content=text, metadata={
                                     "source": str(path), "filetype": "json"}))
                else:
                    print(f"⚠️  Empty JSON: {path.name}")

            # ---- YAML / YML ----
            elif ext in (".yaml", ".yml"):
                text = _yaml_to_text(path)
                if text.strip():
                    documents.append(Document(page_content=text, metadata={
                                     "source": str(path), "filetype": "yaml"}))
                else:
                    print(f"⚠️  Empty YAML: {path.name}")

            # ---- XML ----
            elif ext == ".xml":
                text = _xml_to_text(path)
                if text.strip():
                    documents.append(Document(page_content=text, metadata={
                                     "source": str(path), "filetype": "xml"}))
                else:
                    print(f"⚠️  Empty XML: {path.name}")

            # ---- RTF ----
            elif ext == ".rtf":
                text = _rtf_to_text(path)
                if text.strip():
                    documents.append(Document(page_content=text, metadata={
                                     "source": str(path), "filetype": "rtf"}))
                else:
                    print(f"⚠️  Empty RTF: {path.name}")

            # ---- ODT ----
            elif ext == ".odt":
                text = _odt_to_text(path)
                if text.strip():
                    documents.append(Document(page_content=text, metadata={
                                     "source": str(path), "filetype": "odt"}))
                else:
                    print(f"⚠️  Empty ODT: {path.name}")

            # ---- Изображения (OCR) ----
            elif ext in (".png", ".jpg", ".jpeg", ".tif", ".tiff", ".bmp"):
                txt = _ocr_image_bytes(_read_bytes(path))
                if txt.strip():
                    documents.append(Document(page_content=txt, metadata={
                                     "source": str(path), "filetype": "image"}))
                else:
                    print(f"⚠️  Image has no OCR text: {path.name}")

            else:
                print(f"⚠️  Unsupported format: {path.name} — skipped")

        except Exception as exc:
            print(f"❌ Error reading {path.name}: {exc}")

    return documents


def chunk_documents(
    docs: List[Document],
    chunk_size: int = 1024,
    chunk_overlap: int = 100,
) -> List[Document]:

    splitter = RecursiveCharacterTextSplitter(
        chunk_size=chunk_size,
        chunk_overlap=chunk_overlap,
        separators=["\n\n", "\n", " ", ""],
    )

    chunks: List[Document] = []
    for doc in docs:
        for i, chunk_text in enumerate(splitter.split_text(doc.page_content)):
            meta = dict(doc.metadata)
            meta.update({"chunk_index": i})  # FIX
            chunks.append(Document(page_content=chunk_text, metadata=meta))

    return chunks
