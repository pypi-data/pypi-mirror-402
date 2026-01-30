from importlib.resources import files
from pathlib import Path
from typing import Callable, Literal
import contextlib
import io

import instructor

# Suppress pymupdf4llm import warning about layout package
with contextlib.redirect_stdout(io.StringIO()):
    import pymupdf4llm
from docx import Document
from litellm import completion
from pptx import Presentation
from pydantic import BaseModel, Field

from moves_cli.models import Section


class SectionProducer:
    def _extract_pdf(
        self, file_path: Path, extraction_type: Literal["transcript", "presentation"]
    ) -> str:
        """Extract text from PDF using PyMuPDF4LLM (optimized for LLM processing).

        Uses markdown conversion which preserves structure better than plain text.
        """
        # PyMuPDF4LLM works with file path (str)
        # Using page_chunks for better structure preservation
        chunks = pymupdf4llm.to_markdown(
            str(file_path),
            page_chunks=True,  # Returns list of dicts, one per page
        )

        match extraction_type:
            case "transcript":
                # Extract all text from all pages, concatenate into one line
                full_text = " ".join(chunk["text"] for chunk in chunks)  # type: ignore
                # Remove extra spaces and newlines for transcript
                result = " ".join(full_text.split())
                return result

            case "presentation":
                # For page-by-page presentation, keep markdown structure
                markdown_sections = []
                for i, chunk in enumerate(chunks):
                    # Clean up markdown text but preserve basic structure
                    page_text = chunk["text"].strip()  # type: ignore
                    cleaned_text = " ".join(page_text.split())
                    markdown_sections.append(f"# Slide Page {i + 1}\n{cleaned_text}")

                return "\n\n".join(markdown_sections)

    def _extract_docx(
        self, file_path: Path, extraction_type: Literal["transcript", "presentation"]
    ) -> str:
        """Extract text from DOCX using python-docx (free, no PyMuPDF Pro needed)."""
        # Read document (python-docx accepts str path)
        doc = Document(str(file_path))

        match extraction_type:
            case "transcript":
                # Extract all text from all paragraphs
                full_text = " ".join(paragraph.text for paragraph in doc.paragraphs)
                result = " ".join(full_text.split())
                return result

            case "presentation":
                # Treat each paragraph as a potential slide
                # This is a heuristic - DOCX doesn't have explicit slides
                markdown_sections = []
                for i, paragraph in enumerate(doc.paragraphs):
                    if paragraph.text.strip():  # Skip empty paragraphs
                        cleaned_text = " ".join(paragraph.text.split())
                        markdown_sections.append(
                            f"# Slide Page {i + 1}\n{cleaned_text}"
                        )

                return (
                    "\n\n".join(markdown_sections)
                    if markdown_sections
                    else "# Slide Page 1\n"
                )

    def _extract_pptx(
        self, file_path: Path, extraction_type: Literal["transcript", "presentation"]
    ) -> str:
        """Extract text from PPTX using python-pptx (free, no PyMuPDF Pro needed)."""
        # Load presentation (python-pptx accepts str path)
        prs = Presentation(str(file_path))

        match extraction_type:
            case "transcript":
                # Extract all text from all slides
                all_text = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            for paragraph in shape.text_frame.paragraphs:  # type: ignore
                                for run in paragraph.runs:
                                    all_text.append(run.text)

                full_text = " ".join(all_text)
                result = " ".join(full_text.split())
                return result

            case "presentation":
                # Extract text slide by slide
                markdown_sections = []
                for i, slide in enumerate(prs.slides):
                    slide_text_parts = []
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            for paragraph in shape.text_frame.paragraphs:  # type: ignore
                                for run in paragraph.runs:
                                    slide_text_parts.append(run.text)

                    slide_text = " ".join(slide_text_parts)
                    cleaned_text = " ".join(slide_text.split())
                    markdown_sections.append(f"# Slide Page {i + 1}\n{cleaned_text}")

                return "\n\n".join(markdown_sections)

    def _extract_txt(
        self, file_path: Path, extraction_type: Literal["transcript", "presentation"]
    ) -> str:
        """Extract text from plain text files."""
        # Read as UTF-8 text
        content = file_path.read_text(encoding="utf-8")

        match extraction_type:
            case "transcript":
                # Single line, no extra spaces
                result = " ".join(content.split())
                return result

            case "presentation":
                # For TXT, treat each line as a potential section
                lines = [line.strip() for line in content.splitlines() if line.strip()]
                markdown_sections = [
                    f"# Slide Page {i + 1}\n{line}" for i, line in enumerate(lines)
                ]
                return (
                    "\n\n".join(markdown_sections)
                    if markdown_sections
                    else "# Slide Page 1\n"
                )

    def _extract_document(
        self, file_path: Path, extraction_type: Literal["transcript", "presentation"]
    ) -> str:
        """Extract text from document.

        Supports: PDF (pymupdf4llm), DOCX (python-docx), PPTX (python-pptx), TXT.
        All formats are free and don't require PyMuPDF Pro.
        """
        suffix = file_path.suffix.lower()

        try:
            match suffix:
                case ".pdf":
                    return self._extract_pdf(file_path, extraction_type)
                case ".docx":
                    return self._extract_docx(file_path, extraction_type)
                case ".pptx":
                    return self._extract_pptx(file_path, extraction_type)
                case ".txt":
                    return self._extract_txt(file_path, extraction_type)
                case _:
                    raise ValueError(
                        f"Unsupported file format: {suffix}. "
                        f"Supported formats: .pdf, .docx, .pptx, .txt"
                    )
        except Exception as e:
            raise RuntimeError(
                f"Document extraction failed for {file_path} ({extraction_type}): {e}"
            ) from e

    def generate_template(self, presentation_path: Path) -> list[Section]:
        """
        Extract slide count from presentation file and generate empty sections.
        Supports: PDF, DOCX, PPTX, TXT (all free, no PyMuPDF Pro needed).
        No LLM call, fully offline.
        """
        suffix = presentation_path.suffix.lower()

        try:
            match suffix:
                case ".pdf":
                    # Use pymupdf4llm to get page chunks
                    chunks = pymupdf4llm.to_markdown(
                        str(presentation_path), page_chunks=True
                    )
                    slide_count = len(chunks)

                case ".docx":
                    # Count non-empty paragraphs as slides
                    doc = Document(str(presentation_path))
                    slide_count = sum(1 for p in doc.paragraphs if p.text.strip())
                    slide_count = max(slide_count, 1)  # At least 1 slide

                case ".pptx":
                    # Count slides directly
                    prs = Presentation(str(presentation_path))
                    slide_count = len(prs.slides)

                case ".txt":
                    # Count non-empty lines as slides
                    content = presentation_path.read_text(encoding="utf-8")
                    lines = [
                        line.strip() for line in content.splitlines() if line.strip()
                    ]
                    slide_count = max(len(lines), 1)  # At least 1 slide

                case _:
                    raise ValueError(
                        f"Unsupported file format: {suffix}. "
                        f"Supported formats: .pdf, .docx, .pptx, .txt"
                    )

            return [
                Section(section_index=i + 1, content="") for i in range(slide_count)
            ]
        except Exception as e:
            raise RuntimeError(
                f"Failed to generate template from {presentation_path}: {e}"
            ) from e

    def estimate_for_files(
        self,
        presentation_path: Path,
        transcript_path: Path,
        llm_model: str,
    ) -> tuple[int, int, float | None]:
        """
        Estimate token count and cost for given files without making LLM call.

        Returns:
            tuple: (slide_count, token_count, estimated_cost_usd or None)
        """
        from litellm import cost_per_token, token_counter

        # Extract data from documents (PDF or TXT)
        presentation_data = self._extract_document(presentation_path, "presentation")
        transcript_data = self._extract_document(transcript_path, "transcript")

        # Count slides
        slide_count = len(presentation_data.split("\n\n"))

        # Build messages for token counting
        system_prompt = (
            files("moves_cli.data").joinpath("llm_instruction.md").read_text()
        )
        messages = [
            {"role": "system", "content": system_prompt},
            {
                "role": "user",
                "content": f"Presentation: {presentation_data}\nTranscript: {transcript_data}",
            },
        ]

        # Count tokens (local, free)
        token_count = token_counter(model=llm_model, messages=messages)

        # Estimate cost (local lookup, free)
        try:
            prompt_cost, _ = cost_per_token(
                model=llm_model, prompt_tokens=token_count, completion_tokens=0
            )
        except Exception:
            prompt_cost = None  # Model pricing not available

        return slide_count, token_count, prompt_cost

    def _call_llm(
        self,
        presentation_data: str,
        transcript_data: str,
        llm_model: str,
        llm_api_key: str,
    ) -> list[str]:
        slide_count = len(presentation_data.split("\n\n"))

        # define output model with schema to reliable extract sections from llm response
        class SectionsOutputModel(BaseModel):
            class SectionItem(BaseModel):
                section_index: int = Field(
                    ...,
                    ge=1,
                    description="Index starting from 1",  # descriptions for llm to understand the schema
                )
                content: str = Field(..., description="Content of the section")

            sections: list[SectionItem] = Field(  # type: ignore
                ...,
                description="List of section items, one for each slide",
                min_items=slide_count,  # must be exact number of slides, min or max.
                max_items=slide_count,
            )

        try:
            import warnings

            # silence harmless pydantic serialization warnings from litellm/instructor
            # see: https://github.com/BerriAI/litellm/issues/11759
            warnings.filterwarnings(
                "ignore",
                message="Pydantic serializer warnings",
                category=UserWarning,
            )

            # hmm, i need to rewrite this system prompt for broader use cases, current one is too restrictive
            system_prompt = (
                files("moves_cli.data").joinpath("llm_instruction.md").read_text()
            )
            # we're pathching the litellm with instructor to use any llm with any schema
            client = instructor.from_litellm(
                completion, mode=instructor.Mode.JSON
            )  # json for better llm output quality, also afaik the instructor retries on failure for these

            response = client.chat.completions.create(
                model=llm_model,
                api_key=llm_api_key,
                messages=[
                    {"role": "system", "content": system_prompt},
                    {
                        "role": "user",
                        "content": f"Presentation: {presentation_data}\nTranscript: {transcript_data}",
                    },
                ],
                response_model=SectionsOutputModel,
                temperature=0.2,  # i've set like this for deterministic results but if i change the prompt i need to increase this
            )
            result = [item.content for item in response.sections]
            return result
        except Exception as e:
            raise RuntimeError(f"LLM call failed: {e}") from e

    def convert_to_markdown(self, sections: list[Section]) -> str:
        """Convert sections to markdown format.

        Format:
            # 1. Slide

            Section content here...

            # 2. Slide

            Another section content...
        """
        lines: list[str] = []
        for section in sections:
            lines.append(f"# {section.section_index}. Slide")
            lines.append("")
            if section.content.strip():
                lines.append(section.content.strip())
                lines.append("")

        return "\n".join(lines)

    def load_from_markdown(self, markdown_content: str) -> list[Section]:
        """Load sections from markdown format.

        Parses `# N. Slide` headings as section indices, content follows until next heading.
        """
        import re

        sections: list[Section] = []
        # Split by heading pattern, keeping the delimiter
        pattern = r"^#\s+(\d+)\.\s*Slide\s*$"
        parts = re.split(pattern, markdown_content, flags=re.MULTILINE)

        # parts[0] is content before first heading (usually empty)
        # parts[1], parts[2] = first index, first content
        # parts[3], parts[4] = second index, second content, etc.

        i = 1
        while i < len(parts) - 1:
            section_index = int(parts[i])
            content = parts[i + 1].strip()
            sections.append(Section(section_index=section_index, content=content))
            i += 2

        return sections

    def generate_sections(
        self,
        presentation_path: Path,
        transcript_path: Path,
        llm_model: str,
        llm_api_key: str,
        callback: Callable[[str], None] | None = None,
    ) -> list[Section]:
        if callback:
            callback("Extracting presentation data...")
        presentation_data = self._extract_document(presentation_path, "presentation")

        if callback:
            callback("Extracting transcript data...")
        transcript_data = self._extract_document(transcript_path, "transcript")

        if callback:
            callback("Calling LLM...")
        section_contents = self._call_llm(
            presentation_data=presentation_data,
            transcript_data=transcript_data,
            llm_model=llm_model,
            llm_api_key=llm_api_key,
        )

        generated_sections: list[Section] = []

        for idx, content in enumerate(section_contents):
            section = Section(
                content=content,
                section_index=idx + 1,
            )
            generated_sections.append(section)

        return generated_sections
