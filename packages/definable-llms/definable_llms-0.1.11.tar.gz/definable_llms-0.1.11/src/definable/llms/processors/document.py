"""Document processing for PDF, DOCX, PPTX files."""

import io
import structlog

from ..base.types import FileInfo
from ..base.exceptions import FileProcessingError
from .base import BaseProcessor, TextChunker


logger = structlog.get_logger()


class PDFProcessor(BaseProcessor):
  """PDF document processor using pypdf."""

  def __init__(self, extract_images: bool = False):
    """Initialize PDF processor.

    Args:
        extract_images: Whether to extract images from PDF
    """
    super().__init__()
    self.extract_images = extract_images

  def can_process(self, file_info: FileInfo) -> bool:
    """Check if file is a PDF."""
    return file_info.filename.lower().endswith(".pdf") or file_info.content_type == "application/pdf"

  async def process(self, file_info: FileInfo) -> FileInfo:
    """Process PDF file and extract text.

    Args:
        file_info: PDF file information

    Returns:
        File info with extracted text
    """
    try:
      import pypdf
    except ImportError:
      raise FileProcessingError(
        file_info.filename,
        "pypdf library not available. Install with: pip install pypdf",
      )

    try:
      # Use content if available, otherwise read from path
      if file_info.content:
        pdf_file: io.BytesIO = io.BytesIO(file_info.content)
      elif file_info.path:
        pdf_file = open(file_info.path, "rb")  # type: ignore
      else:
        raise FileProcessingError(file_info.filename, "No content or path provided")

      reader = pypdf.PdfReader(pdf_file)

      # Extract metadata
      metadata = {
        "pages": len(reader.pages),
        "title": reader.metadata.title if reader.metadata else None,
        "author": reader.metadata.author if reader.metadata else None,
        "creator": reader.metadata.creator if reader.metadata else None,
        "producer": reader.metadata.producer if reader.metadata else None,
        "creation_date": str(reader.metadata.creation_date) if reader.metadata and reader.metadata.creation_date else None,
        "modification_date": str(reader.metadata.modification_date) if reader.metadata and reader.metadata.modification_date else None,
      }

      # Extract text from all pages
      text_content = []
      for page_num, page in enumerate(reader.pages):
        try:
          page_text = page.extract_text()
          if page_text.strip():
            text_content.append(f"[Page {page_num + 1}]\n{page_text}")
        except Exception as e:
          self.logger.warning(f"Failed to extract text from page {page_num + 1}: {e}")

      # Close file if we opened it
      if file_info.path and hasattr(pdf_file, "close"):
        pdf_file.close()

      # Join all text content
      processed_text = "\n\n".join(text_content)

      # Update file info
      file_info.processed_text = processed_text
      file_info.metadata.update(metadata)

      # Create chunks
      chunker = TextChunker()
      file_info.chunks = chunker.chunk_text(processed_text)

      self.logger.info(f"Processed PDF: {len(text_content)} pages, {len(processed_text)} characters, {len(file_info.chunks)} chunks")

      return file_info

    except Exception as e:
      self.logger.error(f"Failed to process PDF {file_info.filename}: {e}")
      raise FileProcessingError(file_info.filename, f"PDF processing failed: {e}")


class DocxProcessor(BaseProcessor):
  """DOCX document processor using python-docx."""

  def can_process(self, file_info: FileInfo) -> bool:
    """Check if file is a DOCX."""
    return (
      file_info.filename.lower().endswith(".docx")
      or file_info.content_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    )

  async def process(self, file_info: FileInfo) -> FileInfo:
    """Process DOCX file and extract text.

    Args:
        file_info: DOCX file information

    Returns:
        File info with extracted text
    """
    try:
      import docx
    except ImportError:
      raise FileProcessingError(
        file_info.filename,
        "python-docx library not available. Install with: pip install python-docx",
      )

    try:
      # Use content if available, otherwise read from path
      if file_info.content:
        doc_file: io.BytesIO = io.BytesIO(file_info.content)
      elif file_info.path:
        doc_file = file_info.path  # type: ignore
      else:
        raise FileProcessingError(file_info.filename, "No content or path provided")

      document = docx.Document(doc_file)

      # Extract metadata
      core_props = document.core_properties
      metadata = {
        "title": core_props.title,
        "author": core_props.author,
        "subject": core_props.subject,
        "keywords": core_props.keywords,
        "category": core_props.category,
        "comments": core_props.comments,
        "created": str(core_props.created) if core_props.created else None,
        "modified": str(core_props.modified) if core_props.modified else None,
        "last_modified_by": core_props.last_modified_by,
        "revision": core_props.revision,
        "paragraphs": len(document.paragraphs),
      }

      # Extract text from paragraphs
      text_content = []
      for paragraph in document.paragraphs:
        if paragraph.text.strip():
          text_content.append(paragraph.text.strip())

      # Extract text from tables
      table_content = []
      for table in document.tables:
        for row in table.rows:
          row_text = []
          for cell in row.cells:
            if cell.text.strip():
              row_text.append(cell.text.strip())
          if row_text:
            table_content.append(" | ".join(row_text))

      # Combine all text
      all_content = []
      if text_content:
        all_content.extend(text_content)
      if table_content:
        all_content.append("\n[Tables]")
        all_content.extend(table_content)

      processed_text = "\n\n".join(all_content)

      # Update file info
      file_info.processed_text = processed_text
      file_info.metadata.update(metadata)

      # Create chunks
      chunker = TextChunker()
      file_info.chunks = chunker.chunk_text(processed_text)

      self.logger.info(
        f"Processed DOCX: {len(text_content)} paragraphs, "
        f"{len(table_content)} table rows, "
        f"{len(processed_text)} characters, {len(file_info.chunks)} chunks"
      )

      return file_info

    except Exception as e:
      self.logger.error(f"Failed to process DOCX {file_info.filename}: {e}")
      raise FileProcessingError(file_info.filename, f"DOCX processing failed: {e}")


class PowerPointProcessor(BaseProcessor):
  """PowerPoint PPTX processor using python-pptx."""

  def can_process(self, file_info: FileInfo) -> bool:
    """Check if file is a PPTX."""
    return (
      file_info.filename.lower().endswith(".pptx")
      or file_info.content_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )

  async def process(self, file_info: FileInfo) -> FileInfo:
    """Process PPTX file and extract text.

    Args:
        file_info: PPTX file information

    Returns:
        File info with extracted text
    """
    try:
      from pptx import Presentation
    except ImportError:
      raise FileProcessingError(
        file_info.filename,
        "python-pptx library not available. Install with: pip install python-pptx",
      )

    try:
      # Use content if available, otherwise read from path
      if file_info.content:
        pptx_file: io.BytesIO = io.BytesIO(file_info.content)
      elif file_info.path:
        pptx_file = file_info.path  # type: ignore
      else:
        raise FileProcessingError(file_info.filename, "No content or path provided")

      presentation = Presentation(pptx_file)

      # Extract metadata
      core_props = presentation.core_properties
      metadata = {
        "title": core_props.title,
        "author": core_props.author,
        "subject": core_props.subject,
        "keywords": core_props.keywords,
        "category": core_props.category,
        "comments": core_props.comments,
        "created": str(core_props.created) if core_props.created else None,
        "modified": str(core_props.modified) if core_props.modified else None,
        "last_modified_by": core_props.last_modified_by,
        "revision": core_props.revision,
        "slides": len(presentation.slides),
      }

      # Extract text from slides
      slide_content = []
      for slide_num, slide in enumerate(presentation.slides):
        slide_text = []

        # Extract text from shapes
        for shape in slide.shapes:
          if hasattr(shape, "text") and shape.text.strip():
            slide_text.append(shape.text.strip())
          elif hasattr(shape, "table"):
            # Extract text from table
            for row in shape.table.rows:
              row_text = []
              for cell in row.cells:
                if cell.text.strip():
                  row_text.append(cell.text.strip())
              if row_text:
                slide_text.append(" | ".join(row_text))

        if slide_text:
          slide_content.append(f"[Slide {slide_num + 1}]\n" + "\n".join(slide_text))

      # Combine all content
      processed_text = "\n\n".join(slide_content)

      # Update file info
      file_info.processed_text = processed_text
      file_info.metadata.update(metadata)

      # Create chunks
      chunker = TextChunker()
      file_info.chunks = chunker.chunk_text(processed_text)

      self.logger.info(f"Processed PPTX: {len(slide_content)} slides, {len(processed_text)} characters, {len(file_info.chunks)} chunks")

      return file_info

    except Exception as e:
      self.logger.error(f"Failed to process PPTX {file_info.filename}: {e}")
      raise FileProcessingError(file_info.filename, f"PPTX processing failed: {e}")


class SpreadsheetProcessor(BaseProcessor):
  """Excel XLSX/CSV processor using openpyxl."""

  def can_process(self, file_info: FileInfo) -> bool:
    """Check if file is a spreadsheet."""
    filename_lower = file_info.filename.lower()
    return filename_lower.endswith((".xlsx", ".csv")) or file_info.content_type in [
      "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
      "text/csv",
    ]

  async def process(self, file_info: FileInfo) -> FileInfo:
    """Process spreadsheet file and extract text.

    Args:
        file_info: Spreadsheet file information

    Returns:
        File info with extracted text
    """
    if file_info.filename.lower().endswith(".csv"):
      return await self._process_csv(file_info)
    else:
      return await self._process_xlsx(file_info)

  async def _process_csv(self, file_info: FileInfo) -> FileInfo:
    """Process CSV file."""
    import csv
    import io

    try:
      # Get content as text
      if file_info.content:
        content = self._safe_decode(file_info.content)
      elif file_info.path:
        with open(file_info.path, "r", encoding="utf-8") as f:
          content = f.read()
      else:
        raise FileProcessingError(file_info.filename, "No content or path provided")

      # Parse CSV
      csv_reader = csv.reader(io.StringIO(content))
      rows = list(csv_reader)

      if not rows:
        processed_text = ""
      else:
        # Format as table
        header = rows[0]
        data_rows = rows[1:]

        # Create formatted output
        lines = []
        lines.append(" | ".join(header))
        lines.append("-" * len(lines[0]))

        for row in data_rows[:100]:  # Limit to first 100 rows
          # Pad row to match header length
          padded_row = row + [""] * (len(header) - len(row))
          lines.append(" | ".join(str(cell) for cell in padded_row[: len(header)]))

        if len(data_rows) > 100:
          lines.append(f"... ({len(data_rows) - 100} more rows)")

        processed_text = "\n".join(lines)

      # Metadata
      metadata = {
        "rows": len(rows),
        "columns": len(rows[0]) if rows else 0,
        "has_header": True,  # Assume first row is header
      }

      # Update file info
      file_info.processed_text = processed_text
      file_info.metadata.update(metadata)

      # Create chunks
      chunker = TextChunker(chunk_size=2000)  # Larger chunks for tables
      file_info.chunks = chunker.chunk_text(processed_text)

      self.logger.info(f"Processed CSV: {len(rows)} rows, {len(rows[0]) if rows else 0} columns, {len(file_info.chunks)} chunks")

      return file_info

    except Exception as e:
      self.logger.error(f"Failed to process CSV {file_info.filename}: {e}")
      raise FileProcessingError(file_info.filename, f"CSV processing failed: {e}")

  async def _process_xlsx(self, file_info: FileInfo) -> FileInfo:
    """Process XLSX file."""
    try:
      import openpyxl
    except ImportError:
      raise FileProcessingError(
        file_info.filename,
        "openpyxl library not available. Install with: pip install openpyxl",
      )

    try:
      # Use content if available, otherwise read from path
      if file_info.content:
        xlsx_file: io.BytesIO = io.BytesIO(file_info.content)
      elif file_info.path:
        xlsx_file = file_info.path  # type: ignore
      else:
        raise FileProcessingError(file_info.filename, "No content or path provided")

      workbook = openpyxl.load_workbook(xlsx_file, data_only=True)

      # Extract metadata
      properties = workbook.properties
      metadata = {
        "title": properties.title,
        "creator": properties.creator,
        "subject": properties.subject,
        "description": properties.description,
        "keywords": properties.keywords,
        "category": properties.category,
        "comments": properties.comments,
        "created": str(properties.created) if properties.created else None,
        "modified": str(properties.modified) if properties.modified else None,
        "last_modified_by": properties.lastModifiedBy,
        "worksheets": len(workbook.worksheets),
        "worksheet_names": [ws.title for ws in workbook.worksheets],
      }

      # Extract data from all worksheets
      sheet_content = []
      for worksheet in workbook.worksheets:
        sheet_data = []
        sheet_data.append(f"[Worksheet: {worksheet.title}]")

        # Get all rows with data
        rows = list(worksheet.iter_rows(values_only=True))
        if not rows:
          continue

        # Format as table
        for row_num, row in enumerate(rows[:100]):  # Limit to first 100 rows
          row_data = [str(cell) if cell is not None else "" for cell in row]
          if any(cell.strip() for cell in row_data):  # Skip empty rows
            sheet_data.append(" | ".join(row_data))

        if len(rows) > 100:
          sheet_data.append(f"... ({len(rows) - 100} more rows)")

        if len(sheet_data) > 1:  # More than just the title
          sheet_content.append("\n".join(sheet_data))

      # Combine all sheets
      processed_text = "\n\n".join(sheet_content)

      # Update file info
      file_info.processed_text = processed_text
      file_info.metadata.update(metadata)

      # Create chunks
      chunker = TextChunker(chunk_size=2000)  # Larger chunks for tables
      file_info.chunks = chunker.chunk_text(processed_text)

      self.logger.info(f"Processed XLSX: {len(workbook.worksheets)} sheets, {len(processed_text)} characters, {len(file_info.chunks)} chunks")

      return file_info

    except Exception as e:
      self.logger.error(f"Failed to process XLSX {file_info.filename}: {e}")
      raise FileProcessingError(file_info.filename, f"XLSX processing failed: {e}")


class PlainTextProcessor(BaseProcessor):
  """Plain text file processor."""

  def can_process(self, file_info: FileInfo) -> bool:
    """Check if file is plain text."""
    filename_lower = file_info.filename.lower()
    text_extensions = [
      ".txt",
      ".md",
      ".json",
      ".xml",
      ".yaml",
      ".yml",
      ".py",
      ".js",
      ".html",
      ".css",
    ]
    text_content_types = [
      "text/plain",
      "text/markdown",
      "application/json",
      "text/xml",
      "text/html",
    ]

    return bool(
      any(filename_lower.endswith(ext) for ext in text_extensions)
      or file_info.content_type in text_content_types
      or (file_info.content_type and file_info.content_type.startswith("text/"))
    )

  async def process(self, file_info: FileInfo) -> FileInfo:
    """Process plain text file.

    Args:
        file_info: Text file information

    Returns:
        File info with processed text
    """
    try:
      # Get content as text
      if file_info.content:
        processed_text = self._safe_decode(file_info.content)
      elif file_info.path:
        with open(file_info.path, "r", encoding="utf-8") as f:
          processed_text = f.read()
      else:
        raise FileProcessingError(file_info.filename, "No content or path provided")

      # Extract basic metadata
      metadata = {
        "lines": len(processed_text.splitlines()),
        "characters": len(processed_text),
        "words": len(processed_text.split()),
        "encoding": "utf-8",  # Simplified assumption
      }

      # Update file info
      file_info.processed_text = processed_text
      file_info.metadata.update(metadata)

      # Create chunks
      chunker = TextChunker()
      file_info.chunks = chunker.chunk_text(processed_text)

      self.logger.info(f"Processed text file: {metadata['lines']} lines, {metadata['characters']} characters, {len(file_info.chunks)} chunks")

      return file_info

    except Exception as e:
      self.logger.error(f"Failed to process text file {file_info.filename}: {e}")
      raise FileProcessingError(file_info.filename, f"Text processing failed: {e}")
