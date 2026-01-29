"""
This module integrates language model processing with sprint report generation and interaction,
utilizing data from GitLab and structuring outputs in various formats including JSON, tables,
and PowerPoint presentations. It primarily uses the langchain library to structure and parse
outputs and the rich library to enhance command-line interface outputs.
"""
import os
from rich import print
from rich.table import Table
from pptx import Presentation

from langchain_core.output_parsers import PydanticOutputParser
from langchain_core.prompts import ChatPromptTemplate
from langchain_core.runnables import RunnablePassthrough
from langchain_community.chat_models import ChatOllama
from .prompts import prompts
from .models import SprintReport, ResponseModel


def prepare_parser_and_prompt(model):
    """
    Prepares a parser and prompt template for use in generating structured outputs based on a specified model.

    Parameters:
        model (BaseModel): A Pydantic model that defines the schema for output data.

    Returns:
        tuple: A tuple containing the configured parser and prompt template.
    """
    parser = PydanticOutputParser(pydantic_object=model)
    prompt_template = ChatPromptTemplate.from_messages([
        ("system",
         "Answer the user query. Wrap the output in `json` tags\n{format_instructions}. The context is {context}"),
        ("human", "{query}"),
    ]).partial(format_instructions=parser.get_format_instructions())

    return parser, prompt_template


def prepare_chain(llm, retriever, model):
    """
    Prepares the entire processing chain from prompt preparation to response parsing.

    Parameters:
        llm (ChatOllama): The language model used for generating responses.
        retriever (any): The retriever instance used to fetch context.
        model (BaseModel): The model specifying output schema.

    Returns:
        Pipeline: A fully configured pipeline ready for invoking with user queries.
    """
    parser, prompt_template = prepare_parser_and_prompt(model)
    return {"context": retriever, "query": RunnablePassthrough()} | prompt_template | llm | parser


def invoke_chain(chain, prompt, max_attempts):
    """
    Invokes the processing chain with retries on failure up to a maximum number of attempts.

    Parameters:
        chain (Pipeline): The processing chain to be invoked.
        prompt (str): The user query to be processed.
        max_attempts (int): Maximum number of retry attempts.

    Returns:
        Any: The result from the processing chain or an error message after max attempts.
    """
    count_attempts = 0
    while count_attempts < max_attempts:
        try:
            response = chain.invoke(prompt)
            return response
        except Exception as e:
            count_attempts += 1
    print("[red]❌ Unable to generate answer[/red]")


def get_report(llm, embeddings, max_attempts, sprint_goals=None):
    """
    Generates a structured report based on sprint goals and activities data.

    Parameters:
        llm (ChatOllama): Language model used for data processing.
        embeddings (any): Data retriever containing sprint activity embeddings.
        max_attempts (int): Maximum number of attempts for generating the report.
        sprint_goals (str, optional): Additional sprint goals to be appended in the report.

    Returns:
        Any: Structured report containing sprint activities and statuses.
    """
    chain = prepare_chain(llm, embeddings, SprintReport)
    prompt = prompts['generate_report']['question']
    if sprint_goals is not None:
        prompt += prompts['generate_report']['additional_info'] + sprint_goals
    report = invoke_chain(chain, prompt, max_attempts)
    return report


def report(llm, embeddings, sprint_goals, max_attempts):
    """
    Prints a sprint report in a rich formatted table.

    Parameters:
        llm (ChatOllama): Language model for processing.
        embeddings (any): Embeddings of sprint data.
        sprint_goals (str): Sprint goals to include in the report.
        max_attempts (int): Maximum number of attempts for generating the report.
    """
    print("[bold]Generating report...[/bold]")
    report = get_report(llm, embeddings, max_attempts, sprint_goals)

    table = Table(title="Sprint Activities Report")

    table.add_column("Title", style="cyan", no_wrap=True)
    table.add_column("Description", style="magenta")
    table.add_column("Status", style="green")

    for activity in report.activities:
        table.add_row(
            activity.title,
            activity.brief_desc_status,
            activity.status
        )

    print("[green]✅ Report generated[/green]")
    print(table)


def pptx(llm, embeddings, sprint_goals, pptx_file_name, max_attempts):
    """
    Generates a PowerPoint presentation from sprint report data.

    Parameters:
        llm (ChatOllama): Language model used for generating the report.
        embeddings (any): Embeddings containing sprint data.
        sprint_goals (str): Sprint goals to include in the report.
        pptx_file_name (str): File name for the generated PowerPoint presentation.
        max_attempts (int): Maximum number of attempts for data retrieval.
    """
    report = get_report(llm, embeddings, max_attempts, sprint_goals)

    prs = Presentation()
    slide_layout = prs.slide_layouts[1]

    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]

    title.text = "Sprint Activities Report"
    content.text = "Activities Summary:\n"

    for activity in report.activities:
        content.text += f"- {activity.title}: {activity.brief_desc_status} (Status: {activity.status})\n"

    curr_dir = os.getcwd()
    prs.save(os.path.join(curr_dir, pptx_file_name))
    print("[green]✅ PowerPoint report generated[/green]")


def chat(llm, embeddings, sprint_goals, max_attempts):
    """
    Conducts a chat session where users can interactively query sprint status.

    Parameters:
        llm (ChatOllama): Language model for processing chat queries.
        embeddings (any): Data retriever for sprint activities.
        sprint_goals (str): Sprint goals to include in the chat session.
        max_attempts (int): Maximum number of attempts for responding to each query.
    """
    is_end = False
    while not is_end:
        user_input = input("You: ")
        if user_input == "exit" or user_input == "quit" or user_input == "q":
            is_end = True
            print("[bold]👋 Goodbye![/bold]")
            break
        chain = prepare_chain(llm, embeddings, ResponseModel)
        response = invoke_chain(
            chain,
            user_input + "\n" + prompts['chat']['question'] + "\n Sprint goals were: " + sprint_goals,
            max_attempts
        )
        print(f"Bot: {response}")


def execute_command(command, embeddings, model, max_tokens, max_attempts, sprint_goals, pptx_file_name):
    """
    Executes the specified command using the provided model and data.

    Parameters:
        command (str): The command to execute ('report', 'pptx', or 'chat').
        embeddings (any): Embeddings containing sprint data.
        model (str): Model identifier for the language model.
        max_tokens (int): Maximum number of tokens the model can handle.
        max_attempts (int): Maximum number of attempts for command execution.
        sprint_goals (str): Sprint goals to be included where applicable.
        pptx_file_name (str): Filename for storing PowerPoint presentations.

    Raises:
        Exception: If an invalid command is specified.
    """
    llm = ChatOllama(model=model, max_tokens=max_tokens)
    if command == "report":
        report(llm, embeddings, sprint_goals, max_attempts)
    elif command == "pptx":
        pptx(llm, embeddings, sprint_goals, pptx_file_name, max_attempts)
    elif command == "chat":
        chat(llm, embeddings, sprint_goals, max_attempts)
    else:
        print("Invalid command")
        raise Exception("Invalid command")
