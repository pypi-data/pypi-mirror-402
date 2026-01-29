"""Conversor de arquivos PPTX para Markdown com exportação de mídias."""

from __future__ import annotations

from pptx import Presentation
from pathlib import Path


def convert_pptx_to_markdown(pptx_path: str, md_path: str, media_dir: str = None) -> str:
    """Converte um PPTX em Markdown, extraindo mídias para um diretório.

    Args:
        pptx_path: Caminho absoluto ou relativo para o arquivo `.pptx`.
        md_path: Caminho de saída para o arquivo Markdown gerado.
        media_dir: Diretório onde imagens e gráficos exportados serão salvos.

    Returns:
        String em formato Markdown com os conteúdos do slide separados por
        delimitadores de seção.
    """
    if media_dir is None:
        media_dir = Path(md_path).parent.name

    Path(media_dir).mkdir(parents=True, exist_ok=True)
    prs = Presentation(pptx_path)
    md_lines = []

    for slide_idx, slide in enumerate(prs.slides, start=1):
        md_lines.append(f"# Slide {slide_idx}")
        md_lines.append("")

        for shape_idx, shape in enumerate(slide.shapes, start=1):

            # Texto dos shapes (incluindo listas, níveis, bold/italic)
            if hasattr(shape, "text_frame") and shape.text_frame:
                for para in shape.text_frame.paragraphs:
                    md = paragraph_to_md_pptx(para)
                    if md:
                        md_lines.append(md)
                md_lines.append("")

            # Tabelas
            if shape.has_table:
                md_lines.append(table_to_md_pptx(shape.table))
                md_lines.append("")

            # Imagens
            if shape.shape_type == 13:
                img_path = extract_pptx_image(shape, media_dir, slide_idx, shape_idx)
                md_lines.append(f"![image]({img_path})")
                md_lines.append("")

            # Gráficos
            if hasattr(shape, "chart"):
                try:
                    chart_img = extract_pptx_chart(shape, media_dir, slide_idx, shape_idx)
                    md_lines.append(f"![chart]({chart_img})")
                    md_lines.append("")
                except:
                    pass

        md_lines.append("\n---\n")

    Path(md_path).write_text("\n".join(md_lines), encoding="utf-8")
    return md_path


def run_to_md_pptx(run):
    """Aplica formatação Markdown a um run de texto do PPTX."""
    text = run.text.replace("\n", " ").strip()

    if run.bold:
        text = f"**{text}**"
    if run.italic:
        text = f"*{text}*"

    return text


def paragraph_to_md_pptx(para):
    """Converte parágrafos do PPTX em Markdown, mantendo hierarquia de listas.

    Args:
        para: Parágrafo oriundo de `shape.text_frame` com formatação de nível.

    Returns:
        Texto em Markdown representando o parágrafo com indentação adequada.
    """
    level = para.level if hasattr(para, "level") else 0
    indent = "  " * level
    text = "".join(run_to_md_pptx(r) for r in para.runs).strip()

    if not text:
        return ""

    if para.level > 0:
        return indent + "- " + text

    return text


def extract_pptx_image(shape, output_folder, slide_num, shape_num):
    """Extrai uma imagem de um shape e a salva no disco.

    Args:
        shape: Shape do slide contendo a imagem binária.
        output_folder: Pasta alvo para salvar a mídia.
        slide_num: Índice do slide atual.
        shape_num: Índice do shape dentro do slide.

    Returns:
        Caminho absoluto da imagem criada.
    """
    img = shape.image
    ext = img.ext or "png"
    img_name = f"slide{slide_num}_obj{shape_num}.{ext}"
    img_path = Path(output_folder) / img_name
    img_path.write_bytes(img.blob)
    return str(img_path)


def extract_pptx_chart(shape, output_folder, slide_num, shape_num):
    """Exporta gráficos embutidos no PPTX como PNG.

    Args:
        shape: Shape contendo o gráfico.
        output_folder: Pasta alvo para o arquivo gerado.
        slide_num: Índice do slide atual.
        shape_num: Índice do gráfico dentro do slide.

    Returns:
        Caminho absoluto do PNG exportado.
    """
    chart_part = shape.chart._chart
    # renderização bruta do gráfico como imagem:
    img_bytes = chart_part.part.blob
    img_name = f"slide{slide_num}_chart{shape_num}.png"
    img_path = Path(output_folder) / img_name
    img_path.write_bytes(img_bytes)
    return str(img_path)


def table_to_md_pptx(table):
    """Converte uma tabela de slide em uma string Markdown.

    Args:
        table: Objeto de tabela do PowerPoint.

    Returns:
        Conteúdo da tabela formatado em Markdown.
    """
    md = []
    header = table.rows[0].cells
    md.append("| " + " | ".join(c.text.strip() for c in header) + " |")
    md.append("| " + " | ".join("---" for _ in header) + " |")

    for row in table.rows[1:]:
        md.append("| " + " | ".join(c.text.strip() for c in row.cells) + " |")

    return "\n".join(md)

