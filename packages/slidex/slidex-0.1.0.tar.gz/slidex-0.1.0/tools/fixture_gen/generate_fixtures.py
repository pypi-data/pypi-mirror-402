import tempfile
from pathlib import Path
import struct
import zlib

from pptx import Presentation
from pptx.util import Inches

ROOT = Path(__file__).resolve().parents[2]
FIXTURES = ROOT / "tests" / "fixtures" / "generated" / "simple"


def png_bytes(width: int, height: int, rgb: tuple[int, int, int]) -> bytes:
    r, g, b = rgb
    raw = bytearray()
    row = bytes([r, g, b]) * width
    for _ in range(height):
        raw.append(0)
        raw.extend(row)
    compressed = zlib.compress(bytes(raw))

    def chunk(tag: bytes, data: bytes) -> bytes:
        return struct.pack(">I", len(data)) + tag + data + struct.pack(">I", zlib.crc32(tag + data) & 0xFFFFFFFF)

    signature = b"\x89PNG\r\n\x1a\n"
    ihdr = struct.pack(">IIBBBBB", width, height, 8, 2, 0, 0, 0)
    return signature + chunk(b"IHDR", ihdr) + chunk(b"IDAT", compressed) + chunk(b"IEND", b"")


def title_and_content() -> None:
    FIXTURES.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(layout)

    slide.shapes.title.text = "Hello World"
    slide.placeholders[1].text = "This is text"

    prs.save(FIXTURES / "title_and_content.pptx")


def text_only() -> None:
    FIXTURES.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    left = top = width = height = 1000000
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.text = "First line"
    p = tf.add_paragraph()
    p.text = "Second line"
    p = tf.add_paragraph()
    p.text = "Hello {{name}}"

    prs.save(FIXTURES / "text_only.pptx")


def blank_template() -> None:
    assets = ROOT / "assets"
    assets.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[6])
    prs.save(assets / "blank.pptx")


def picture_simple() -> None:
    FIXTURES.mkdir(parents=True, exist_ok=True)

    with tempfile.TemporaryDirectory() as tmpdir:
        image_path = Path(tmpdir) / "image.png"
        image_path.write_bytes(png_bytes(64, 64, (220, 0, 0)))

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.add_picture(str(image_path), Inches(1), Inches(1), width=Inches(2), height=Inches(2))
        prs.save(FIXTURES / "picture_simple.pptx")


def main() -> None:
    title_and_content()
    text_only()
    blank_template()
    picture_simple()


if __name__ == "__main__":
    main()
