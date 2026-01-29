# file: algosystem/backtesting/dashboard/utils/ip_slide_generator.py
"""
Enhanced IP Slide Generator with Professional Plotly Chart Styling
- Valid Plotly axis API.
- Higher-res PNG (scale=3) and 1.25x width (1500x600).
- One combined drawdown chart (strategy + benchmark).
- Remove drawdown duration plot.
- Add Rolling Calmar + Rolling Var to risk charts/exports.
- Bigger text (2x) + legend inside for strategy vs benchmark charts.
"""

from __future__ import annotations

import os
import re
import glob
from datetime import datetime
from typing import Optional, Tuple, List

import pandas as pd
import numpy as np

import plotly.graph_objects as go
from pptx import Presentation
from pptx.util import Inches

# ===== Plotly visual defaults =====
_PLOTLY_TEMPLATE = "plotly_white"
_COLORWAY = [
    "#1f77b4", "#ff7f0e", "#2ca02c", "#d62728", "#9467bd",
    "#8c564b", "#e377c2", "#7f7f7f", "#bcbd22", "#17becf",
]
_FONT_FAMILY = "Inter, Segoe UI, Arial, sans-serif"
_TITLE_COLOR = "#2c3e50"
_AXIS_COLOR = "#34495e"
_GRID_COLOR = "#e5e7eb"
_BG_COLOR = "white"

# Size & resolution
_WIDTH_MULTIPLIER = 1.25
_BASE_W, _BASE_H = 1200, 600
_FIG_W, _FIG_H = int(_BASE_W * _WIDTH_MULTIPLIER), _BASE_H  # 1500 x 600
_EXPORT_SCALE = 3  # higher DPI via kaleido scaling

# --- Text scaling (2x = +100%) ---
_TEXT_SCALE = 2.0
_TITLE_SIZE = int(20 * _TEXT_SCALE)
_AXIS_TITLE_SIZE = int(12 * _TEXT_SCALE)
_TICK_SIZE = int(11 * _TEXT_SCALE)
_LEGEND_FONT_SIZE = int(12 * _TEXT_SCALE)

# ---------- Utils ----------

def _slug(name: str) -> str:
    """Filesystem-safe slug."""
    return re.sub(r"[^a-zA-Z0-9]+", "_", (name or "").strip().lower()).strip("_")


def _pick_latest(pattern: str) -> Optional[str]:
    matches = glob.glob(pattern)
    if not matches:
        return None
    def key_fn(p):
        m = re.search(r'_(\d{8}_\d{6})\.', os.path.basename(p))
        if m:
            try:
                return datetime.strptime(m.group(1), "%Y%m%d_%H%M%S")
            except ValueError:
                pass
        return datetime.fromtimestamp(os.path.getmtime(p))
    return sorted(matches, key=key_fn)[-1]


def _load_timeseries_csv(path: str) -> pd.DataFrame:
    """Load CSV and ensure index is datetime."""
    df = pd.read_csv(path)

    # Try to find and convert date column
    date_col = None
    for col in df.columns:
        if col.lower() in ['date', 'datetime', 'timestamp', '']:
            date_col = col
            break

    # If no explicit date column found, try first column
    if date_col is None:
        date_col = df.columns[0]

    try:
        # Attempt datetime conversion (removed deprecated infer_datetime_format)
        df[date_col] = pd.to_datetime(df[date_col], errors='coerce')

        # Check if conversion was successful (most values are datetime, not NaT)
        non_nat_ratio = df[date_col].notna().sum() / len(df)
        if non_nat_ratio > 0.9:  # At least 90% successfully converted
            df = df.set_index(date_col).sort_index()
            # Ensure index is datetime64
            if df.index.dtype != 'datetime64[ns]':
                df.index = pd.to_datetime(df.index, errors='coerce')
        # Note: If conversion fails, we return original DataFrame unchanged
    except Exception:
        # Conversion failed, return DataFrame unchanged
        pass

    return df


def _ensure_dir(d: str) -> None:
    os.makedirs(d, exist_ok=True)


def _pct_like(colname: str) -> bool:
    return any(k in colname.lower() for k in [
        "return", "drawdown", "relative", "vol", "sharpe", "sortino", "var", "calmar"
    ])


def _ensure_datetime_index(series_or_df) -> pd.Series | pd.DataFrame:
    """Ensure the index is datetime, but only if it looks like dates."""
    index_dtype = series_or_df.index.dtype

    # If already datetime, return unchanged
    if index_dtype == 'datetime64[ns]':
        return series_or_df

    # If not datetime, try to convert only if it looks like datetime strings
    try:
        # Check if index values look like dates (first few values)
        sample = str(series_or_df.index[0]).strip()
        # Simple heuristic: dates usually have dashes or slashes
        if '-' in sample or '/' in sample or len(sample) >= 8:
            # Attempt conversion with coerce to avoid NaT explosion
            new_index = pd.to_datetime(series_or_df.index, errors='coerce')
            # Only use if conversion was mostly successful
            if new_index.notna().sum() / len(new_index) > 0.9:
                series_or_df.index = new_index
    except Exception:
        pass

    return series_or_df


def _convert_index_to_iso_strings(series_or_df) -> pd.Series | pd.DataFrame:
    """Convert datetime index to ISO format strings for Plotly compatibility."""
    if series_or_df.index.dtype == 'datetime64[ns]':
        # Convert DatetimeIndex to ISO format strings (YYYY-MM-DD)
        series_or_df.index = series_or_df.index.strftime('%Y-%m-%d')
    return series_or_df

# ---------- Plotly core ----------

def _make_fig_base(title: str, y_is_pct: bool) -> go.Figure:
    fig = go.Figure()
    fig.update_layout(
        template=_PLOTLY_TEMPLATE,
        colorway=_COLORWAY,
        title=dict(
            text=title, x=0.03, xanchor="left",
            font=dict(size=_TITLE_SIZE, family=_FONT_FAMILY, color=_TITLE_COLOR),
        ),
        margin=dict(l=60, r=30, t=60, b=60),
        hovermode="x unified",
        legend=dict(
            orientation="h", yanchor="bottom", y=1.02, xanchor="left", x=0.0,
            bgcolor="rgba(255,255,255,0.8)", bordercolor="#d1d5db", borderwidth=1,
            font=dict(size=_LEGEND_FONT_SIZE, family=_FONT_FAMILY),
            itemclick="toggleothers", itemdoubleclick="toggle",
        ),
        plot_bgcolor=_BG_COLOR, paper_bgcolor=_BG_COLOR,
        width=_FIG_W, height=_FIG_H,
    )
    fig.update_xaxes(
        title=dict(text="Date", font=dict(size=_AXIS_TITLE_SIZE, color=_AXIS_COLOR, family=_FONT_FAMILY)),
        tickfont=dict(size=_TICK_SIZE, color=_AXIS_COLOR, family=_FONT_FAMILY),
        type="date",  # Explicitly set as date type for proper formatting
        tickformat="%Y-%m-%d",
        showgrid=True, gridcolor=_GRID_COLOR, linecolor="#cbd5e1", zeroline=False,
        rangeslider=dict(visible=False),
    )
    fig.update_yaxes(
        title=dict(
            text="Value (%)" if y_is_pct else "Value",
            font=dict(size=_AXIS_TITLE_SIZE, color=_AXIS_COLOR, family=_FONT_FAMILY),
        ),
        tickformat=".0%" if y_is_pct else ",",
        tickfont=dict(size=_TICK_SIZE, color=_AXIS_COLOR, family=_FONT_FAMILY),
        showgrid=True, gridcolor=_GRID_COLOR, linecolor="#cbd5e1", zeroline=False,
    )
    return fig


def _fig_to_png(fig: go.Figure, outfile: str) -> str:
    _ensure_dir(os.path.dirname(outfile) or ".")
    try:
        fig.write_image(outfile, scale=_EXPORT_SCALE, engine="kaleido")
    except ValueError as e:
        # why: needs kaleido installed for static PNGs
        raise RuntimeError(
            "Plotly static image export requires 'kaleido'. Install: pip install -U kaleido"
        ) from e
    return outfile

# ---------- Chart builders ----------

def _save_line_chart(
    series_or_df: pd.Series | pd.DataFrame,
    title: str,
    y_as_percent: bool = False,
    outfile: str = "chart.png",
) -> str:
    # Ensure datetime index
    series_or_df = _ensure_datetime_index(series_or_df)

    if "drawdown" in title.lower():
        return _save_drawdown_chart(series_or_df, title, outfile)
    if isinstance(series_or_df, pd.DataFrame) and len(series_or_df.columns) > 1:
        if "equity" in title.lower():
            return _save_comparison_chart(series_or_df, title, y_as_percent, outfile)

    # Convert datetime index to ISO strings for Plotly
    series_or_df = _convert_index_to_iso_strings(series_or_df)

    fig = _make_fig_base(title, y_as_percent)
    if isinstance(series_or_df, pd.Series):
        name = series_or_df.name or "Series"
        fill = "tozeroy" if any(k in title.lower() for k in ["equity", "value"]) else None
        fig.add_trace(go.Scatter(
            x=series_or_df.index, y=series_or_df.values,
            name=name, mode="lines", line=dict(width=3), fill=fill,
            hovertemplate="%{x}<br>%{y}<extra>"+name+"</extra>",
        ))
    else:
        for col in series_or_df.columns:
            fig.add_trace(go.Scatter(
                x=series_or_df.index, y=series_or_df[col],
                name=col, mode="lines", line=dict(width=2.5),
                hovertemplate="%{x}<br>%{y}<extra>"+col+"</extra>",
            ))
    return _fig_to_png(fig, outfile)


def _save_comparison_chart(
    df: pd.DataFrame, title: str, y_as_percent: bool = False, outfile: str = "chart.png",
) -> str:
    # Ensure datetime index
    df = _ensure_datetime_index(df)
    # Convert datetime index to ISO strings for Plotly
    df = _convert_index_to_iso_strings(df)
    fig = _make_fig_base(title, y_as_percent)
    for col in df.columns:
        is_benchmark = "benchmark" in col.lower()
        fig.add_trace(go.Scatter(
            x=df.index, y=df[col], name=col, mode="lines",
            line=dict(width=3 if not is_benchmark else 2.5,
                      dash="solid" if not is_benchmark else "dash"),
            hovertemplate="%{x}<br>%{y}<extra>"+col+"</extra>",
        ))
    # Place legend inside for strategy vs benchmark charts
    fig.update_layout(
        legend=dict(
            x=0.02, y=0.98, xanchor="left", yanchor="top",
            orientation="v",
            bgcolor="rgba(255,255,255,0.8)",
            bordercolor="#d1d5db", borderwidth=1,
            font=dict(size=_LEGEND_FONT_SIZE, family=_FONT_FAMILY),
        )
    )
    return _fig_to_png(fig, outfile)


def _save_drawdown_chart(
    series_or_df: pd.Series | pd.DataFrame, title: str, outfile: str = "chart.png",
) -> str:
    # Ensure datetime index
    series_or_df = _ensure_datetime_index(series_or_df)
    # Convert datetime index to ISO strings for Plotly
    series_or_df = _convert_index_to_iso_strings(series_or_df)
    fig = _make_fig_base(title, y_is_pct=True)
    fig.update_yaxes(title=dict(text="Drawdown (%)",
                                font=dict(size=_AXIS_TITLE_SIZE, color=_AXIS_COLOR, family=_FONT_FAMILY)),
                     tickformat=".0%")

    multi_series = False
    if isinstance(series_or_df, pd.Series):
        name = series_or_df.name or "Drawdown"
        fig.add_trace(go.Scatter(
            x=series_or_df.index, y=series_or_df.values, name=name,
            mode="lines", line=dict(width=2.5), fill="tozeroy",
            hovertemplate="%{x}<br>%{y:.1%}<extra>"+name+"</extra>",
        ))
    else:
        cols = list(series_or_df.columns)
        multi_series = len(cols) > 1
        for col in cols:
            is_benchmark = "benchmark" in col.lower()
            fig.add_trace(go.Scatter(
                x=series_or_df.index, y=series_or_df[col], name=col,
                mode="lines",
                line=dict(width=3 if not is_benchmark else 2.5,
                          dash="solid" if not is_benchmark else "dash"),
                fill="tozeroy",
                hovertemplate="%{x}<br>%{y:.1%}<extra>"+col+"</extra>",
            ))
    fig.add_hline(y=0, line_width=2, line_color="#111827", opacity=0.9)

    try:
        arr = series_or_df.values if isinstance(series_or_df, pd.Series) else series_or_df.values.flatten()
        if len(arr):
            ymin, ymax = float(np.nanmin(arr)), float(np.nanmax(arr))
            pad = 0.05
            fig.update_yaxes(range=[min(ymin - pad, -0.1), max(ymax + pad, 0.05)])
    except Exception:
        pass

    # Legend inside when comparing strategy vs benchmark
    if multi_series:
        fig.update_layout(
            legend=dict(
                x=0.02, y=0.98, xanchor="left", yanchor="top",
                orientation="v",
                bgcolor="rgba(255,255,255,0.8)",
                bordercolor="#d1d5db", borderwidth=1,
                font=dict(size=_LEGEND_FONT_SIZE, family=_FONT_FAMILY),
            )
        )

    return _fig_to_png(fig, outfile)

# ---------- Public API ----------

def create_charts_from_exports(
    output_dir: str = "backtest_exports",
    prefix: str = "backtest",
    charts_dir: str = "charts",
    pptx_template: Optional[str] = None,
    pptx_out: Optional[str] = None,
):
    _ensure_dir(charts_dir)
    created: List[Tuple[str, str]] = []

    # 1) Combined timeseries
    ts_path = _pick_latest(os.path.join(output_dir, f"{prefix}_timeseries_*.csv"))
    if ts_path:
        ts_df = _load_timeseries_csv(ts_path)
        title = "Backtest Time Series"
        outfile = os.path.join(charts_dir, "timeseries_all.png")
        numeric_cols = ts_df.select_dtypes("number").columns
        y_as_percent = len(numeric_cols) > 0 and all(_pct_like(c) for c in numeric_cols)
        if len(numeric_cols) > 0:
            created.append((title, _save_line_chart(ts_df[numeric_cols], title, y_as_percent, outfile)))
            for col in numeric_cols:
                title_c = f"Time Series — {col}"
                outfile_c = os.path.join(charts_dir, f"timeseries_{_slug(col)}.png")
                created.append((title_c, _save_line_chart(ts_df[col].dropna(), title_c, _pct_like(col), outfile_c)))

    # 2) Equity curves + relative
    eq_path = _pick_latest(os.path.join(output_dir, f"{prefix}_equity_curve_*.csv"))
    if eq_path:
        eq_df = _load_timeseries_csv(eq_path)
        cols = [c for c in eq_df.columns if "equity" in c.lower()]
        if cols:
            title = "Equity Curves"
            outfile = os.path.join(charts_dir, "equity_curves.png")
            created.append((title, _save_line_chart(eq_df[cols], title, False, outfile)))
        rel_cols = [c for c in eq_df.columns if "relative" in c.lower()]
        for col in rel_cols:
            title_r = f"Relative Performance — {col}"
            outfile_r = os.path.join(charts_dir, f"relative_{_slug(col)}.png")
            created.append((title_r, _save_line_chart(eq_df[col].dropna(), title_r, True, outfile_r)))

    # 3) Drawdown (ONE combined chart: strategy + benchmark)
    dd_path = _pick_latest(os.path.join(output_dir, f"{prefix}_drawdown_*.csv"))
    if dd_path:
        dd_df = _load_timeseries_csv(dd_path)
        # pick drawdown columns, exclude any 'duration'
        dd_cols = [c for c in dd_df.columns if "drawdown" in c.lower() and "duration" not in c.lower()]
        if dd_cols:
            # prefer benchmark vs non-benchmark
            bench_cols = [c for c in dd_cols if "benchmark" in c.lower()]
            strat_cols = [c for c in dd_cols if "benchmark" not in c.lower()]
            chosen = []
            if strat_cols:
                chosen.append(strat_cols[0])
            if bench_cols:
                chosen.append(bench_cols[0])
            else:
                # fallback to second series if available
                if len(dd_cols) > 1:
                    chosen.append(dd_cols[1])
            dd_plot = dd_df[chosen].dropna(how="all")
            title_d = "Drawdown — Strategy vs Benchmark"
            outfile_d = os.path.join(charts_dir, "drawdown_combined.png")
            created.append((title_d, _save_drawdown_chart(dd_plot, title_d, outfile_d)))

    # 4) Risk metrics over time (add Rolling Calmar & Rolling Var; remove duration)
    risk_path = _pick_latest(os.path.join(output_dir, f"{prefix}_risk_metrics_*.csv"))
    if risk_path:
        risk_df = _load_timeseries_csv(risk_path)
        for col in risk_df.columns:
            if "duration" in col.lower():  # skip drawdown duration
                continue
            title_rm = f"Risk Metric — {col}"
            outfile_rm = os.path.join(charts_dir, f"risk_{_slug(col)}.png")
            created.append((title_rm, _save_line_chart(risk_df[col].dropna(), title_rm, _pct_like(col), outfile_rm)))

    # 5) Raw portfolio value
    pv_path = _pick_latest(os.path.join(output_dir, f"{prefix}_portfolio_value_*.csv"))
    if pv_path:
        pv_df = _load_timeseries_csv(pv_path)
        value_col = None
        for c in pv_df.columns:
            if "value" in c.lower() or "equity" in c.lower() or "portfolio" in c.lower():
                value_col = c
                break
        if value_col:
            title_pv = "Portfolio Value"
            outfile_pv = os.path.join(charts_dir, "portfolio_value.png")
            created.append((title_pv, _save_line_chart(pv_df[value_col].dropna(), title_pv, False, outfile_pv)))

    # Optional PPTX (unchanged)
    if pptx_out:
        if pptx_template and os.path.exists(pptx_template):
            prs = Presentation(pptx_template)
        else:
            prs = Presentation()
            title_only_layout = None
            for i, layout in enumerate(prs.slide_layouts):
                if any(getattr(ph.placeholder_format, "type", None) == 1 for ph in layout.placeholders):
                    title_only_layout = layout
                    break
            if title_only_layout is None:
                title_only_layout = prs.slide_layouts[0]

        for title, img_path in created:
            slide_layout = None
            for layout in prs.slide_layouts:
                if any(getattr(ph.placeholder_format, "type", None) == 1 for ph in layout.placeholders):
                    slide_layout = layout
                    break
            if slide_layout is None:
                slide_layout = prs.slide_layouts[0]

            slide = prs.slides.add_slide(slide_layout)
            if slide.shapes.title:
                slide.shapes.title.text = title

            pic_left = Inches(0.7)
            pic_top = Inches(1.6)
            pic_height = Inches(5.0)
            slide.shapes.add_picture(img_path, pic_left, pic_top, height=pic_height)

        prs.save(pptx_out)

    return created


def export_backtest_to_csv(results, output_dir="backtest_exports", prefix="backtest"):
    os.makedirs(output_dir, exist_ok=True)
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    exported_files = {}

    summary_data = {
        'Metric': ['Initial Capital', 'Final Capital', 'Total Return', 'Start Date', 'End Date'],
        'Value': [
            results.get('initial_capital', 'N/A'),
            results.get('final_capital', 'N/A'),
            f"{results.get('returns', 0) * 100:.2f}%",
            results.get('start_date', 'N/A'),
            results.get('end_date', 'N/A')
        ]
    }

    metrics_data = results.get('metrics', {})
    # Metrics that should be displayed as percentages (multiply by 100)
    percentage_metrics = [
        'total_return', 'annualized_return', 'annual_return', 'max_drawdown',
        'volatility', 'annualized_volatility', 'monthly_volatility',
        'var_95', 'cvar_95', 'alpha', 'tracking_error',
        'best_month', 'worst_month', 'avg_monthly_return',
        'pct_positive_days', 'pct_positive_months',
        'capture_ratio_up', 'capture_ratio_down'
    ]
    # Metrics that should be displayed as plain values (no percentage)
    ratio_metrics = ['sharpe_ratio', 'sortino_ratio', 'calmar_ratio', 'beta',
                     'correlation', 'information_ratio', 'skewness']

    for key, value in metrics_data.items():
        if not key.endswith('_error'):
            metric_name = key.replace('_', ' ').title()
            if isinstance(value, (int, float)):
                if key in percentage_metrics:
                    formatted_value = f"{value * 100:.2f}%"
                elif key in ratio_metrics:
                    formatted_value = f"{value:.2f}"
                elif isinstance(value, int) or value == int(value):
                    formatted_value = f"{int(value)}"
                else:
                    formatted_value = f"{value:.4f}"
            else:
                formatted_value = str(value)
            summary_data['Metric'].append(metric_name)
            summary_data['Value'].append(formatted_value)

    summary_df = pd.DataFrame(summary_data)
    summary_file = os.path.join(output_dir, f"{prefix}_summary_{timestamp}.csv")
    summary_df.to_csv(summary_file, index=False)
    exported_files['summary'] = summary_file

    plots_data = results.get('plots', {})
    time_series_dfs = {}
    for key, series_data in plots_data.items():
        if isinstance(series_data, pd.Series):
            col_name = key.replace('_', ' ').title()
            # Remove duplicate indices, keeping the first occurrence
            series_data = series_data[~series_data.index.duplicated(keep='first')]
            time_series_dfs[col_name] = series_data

    if time_series_dfs:
        combined_ts = pd.concat(time_series_dfs, axis=1)
        combined_ts.index.name = 'Date'

        timeseries_file = os.path.join(output_dir, f"{prefix}_timeseries_{timestamp}.csv")
        combined_ts.to_csv(timeseries_file, date_format="%Y-%m-%d")
        exported_files['timeseries'] = timeseries_file

        _export_chart_specific_data(combined_ts, output_dir, prefix, timestamp, exported_files)

    if 'equity' in results and isinstance(results['equity'], pd.Series):
        equity_series = pd.DataFrame({
            'Date': results['equity'].index,
            'Portfolio Value': results['equity'].values
        })
        raw_equity_file = os.path.join(output_dir, f"{prefix}_portfolio_value_{timestamp}.csv")
        equity_series.to_csv(raw_equity_file, index=False, date_format="%Y-%m-%d")
        exported_files['portfolio_value'] = raw_equity_file

    _create_metadata_file(results, exported_files, output_dir, prefix, timestamp)

    print("\nGenerating charts...")
    charts = create_charts_from_exports(
        output_dir=output_dir, prefix=prefix, charts_dir="charts",
        pptx_template=None, pptx_out=None
    )
    print(f"Created {len(charts)} charts in 'charts/' directory")

    from algosystem.backtesting.dashboard.utils.create_pptx import create_backtest_pptx

    base_dir = os.path.dirname(os.path.abspath(__file__))
    template_path = os.path.join(base_dir, "template.pptx")

    create_backtest_pptx(
        template_path=template_path,
        output_dir="backtest_exports",
        charts_dir="charts",
        output_path="backtest_presentation.pptx",
        prefix="backtest"
    )

    print(f"Created results slides")

    return exported_files


def _export_chart_specific_data(combined_ts, output_dir, prefix, timestamp, exported_files):
    # Equity-related
    equity_data = pd.DataFrame()
    for col in ['Equity Curve', 'Benchmark Equity Curve', 'Relative Performance']:
        if col in combined_ts.columns:
            equity_data[col] = combined_ts[col]
    if not equity_data.empty:
        equity_file = os.path.join(output_dir, f"{prefix}_equity_curve_{timestamp}.csv")
        equity_data.to_csv(equity_file, date_format="%Y-%m-%d")
        exported_files['equity_curve'] = equity_file

    # Drawdown (exclude duration)
    drawdown_data = pd.DataFrame()
    for col in ['Drawdown Series', 'Benchmark Drawdown Series']:
        if col in combined_ts.columns:
            drawdown_data[col] = combined_ts[col]
    if not drawdown_data.empty:
        drawdown_file = os.path.join(output_dir, f"{prefix}_drawdown_{timestamp}.csv")
        drawdown_data.to_csv(drawdown_file, date_format="%Y-%m-%d")
        exported_files['drawdown'] = drawdown_file

    # Risk Metrics Over Time (include Rolling Calmar & Rolling Var)
    risk_data = pd.DataFrame()
    risk_cols = ['Rolling Sharpe', 'Rolling Sortino', 'Rolling Volatility',
                 'Rolling Skew', 'Rolling Var', 'Rolling Calmar']
    for col in risk_cols:
        if col in combined_ts.columns:
            risk_data[col] = combined_ts[col]
    if not risk_data.empty:
        risk_file = os.path.join(output_dir, f"{prefix}_risk_metrics_{timestamp}.csv")
        risk_data.to_csv(risk_file, date_format="%Y-%m-%d")
        exported_files['risk_metrics'] = risk_file


def _create_metadata_file(results, exported_files, output_dir, prefix, timestamp):
    metadata = {
        'Export Timestamp': timestamp,
        'Files Exported': list(exported_files.keys()),
        'Start Date': results.get('start_date', 'N/A'),
        'End Date': results.get('end_date', 'N/A'),
        'Initial Capital': results.get('initial_capital', 'N/A'),
        'Final Capital': results.get('final_capital', 'N/A'),
        'Total Return': f"{results.get('returns', 0) * 100:.2f}%"
    }
    metadata_file = os.path.join(output_dir, f"{prefix}_metadata_{timestamp}.txt")
    with open(metadata_file, 'w') as f:
        for key, value in metadata.items():
            if key == 'Files Exported':
                f.write(f"{key}:\n")
                for file in value:
                    f.write(f"  - {file}\n")
            else:
                f.write(f"{key}: {value}\n")
    exported_files['metadata'] = metadata_file


if __name__ == "__main__":
    import argparse
    parser = argparse.ArgumentParser(description='Generate professional Plotly charts from backtest exports')
    parser.add_argument('--output_dir', default='backtest_exports', help='Directory with CSV exports')
    parser.add_argument('--charts_dir', default='charts', help='Output directory for charts')
    parser.add_argument('--prefix', default='backtest', help='File prefix')
    parser.add_argument('--template', default=None, help='PowerPoint template path')
    parser.add_argument('--pptx_out', default=None, help='Output PowerPoint file')
    args = parser.parse_args()

    charts = create_charts_from_exports(
        output_dir=args.output_dir, prefix=args.prefix,
        charts_dir=args.charts_dir, pptx_template=args.template, pptx_out=args.pptx_out
    )

    print(f"\nCreated {len(charts)} charts:")
    for title, path in charts:
        print(f"  - {title}: {path}")
