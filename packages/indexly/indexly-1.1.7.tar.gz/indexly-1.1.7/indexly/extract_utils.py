"""
📄 extract_utils.py

Purpose:
    Contains file-type specific text extractors (DOCX, EML, MSG).

Key Features:
    - _extract_docx(): Parses Word documents and extracts text and tables.
    - _extract_eml(): Parses .eml email files using eml_parser.
    - _extract_msg(): Parses .msg Outlook files using extract_msg.
    - Each extractor also supports virtual tag detection via fts_core.

Usage:
    Called by `filetype_utils.py -> extract_text_from_file()` during indexing.
"""

# --- stdlib (safe) ---
import io
import re
import os
import sys
import struct
import json
import sqlite3
import subprocess
import shutil
import platform
from datetime import datetime
from contextlib import suppress

# --- package check helpers (migrated from utils.py) ---
def prompt_install(package_list):
    install_all = False
    try:
        for module, package in package_list:
            try:
                __import__(module)
            except ImportError:
                if not install_all:
                    response = (
                        input(f"Install missing package '{package}'? [Y/n/A=all]: ")
                        .strip()
                        .lower()
                    )
                    if response in ("a", "all"):
                        install_all = True
                        response = "y"

                if install_all or response in ("", "y", "yes"):
                    subprocess.check_call(
                        [sys.executable, "-m", "pip", "install", package]
                    )
    except KeyboardInterrupt:
        print("\n⛔ Package installation cancelled by user.")
        sys.exit(1)


def check_and_install_packages(pkg_list):
    try:
        prompt_install(pkg_list)
    except KeyboardInterrupt:
        print("\n❌ Cancelled while checking packages.")
        sys.exit(1)


# ---------------------------------------------------------------------
# External tool checks (boolean-returning for doctor / programmatic use)
# ---------------------------------------------------------------------
def check_exiftool_available() -> bool:
    """Return True if ExifTool is in PATH, else False"""
    return shutil.which("exiftool") is not None


def check_tesseract_available() -> bool:
    """Return True if Tesseract OCR is in PATH, else False"""
    return shutil.which("tesseract") is not None


def print_external_tools_info():
    """Existing print-based behavior (optional, for CLI)"""
    if not check_exiftool_available():
        print("⚠️ ExifTool not found. Install: https://exiftool.org/")

    if not check_tesseract_available():
        os_name = platform.system().lower()
        print("⚠️ Tesseract OCR not found. Install:")
        if "windows" in os_name:
            print("  choco install tesseract OR winget install tesseract")
        elif "darwin" in os_name:
            print("  brew install tesseract")
        elif "linux" in os_name:
            print("  sudo apt install tesseract-ocr")


# --- run checks BEFORE heavy imports ---
check_and_install_packages(
    [
        ("nltk", "nltk"),
        ("fitz", "pymupdf"),
        ("pytesseract", "pytesseract"),
        ("PIL", "pillow"),
        ("docx", "python-docx"),
        ("openpyxl", "openpyxl"),
        ("rapidfuzz", "rapidfuzz"),
        ("fpdf", "fpdf2"),
        ("reportlab", "reportlab"),
        ("bs4", "beautifulsoup4"),
        ("extract_msg", "extract_msg"),
        ("eml_parser", "eml-parser"),
        ("PyPDF2", "PyPDF2"),
        ("watchdog", "watchdog"),
        ("colorama", "colorama"),
        ("pptx", "python-pptx"),
        ("ebooklib", "ebooklib"),
        ("odf", "odfpy"),
        ("pandas", "pandas"),
    ]
)


# --- third-party imports (safe now) ---
import docx
import extract_msg
import eml_parser
import fitz  # PyMuPDF
import pytesseract
import openpyxl
import pandas as pd

from pptx import Presentation
from ebooklib import epub
from bs4 import BeautifulSoup
from odf.opendocument import load
from odf.text import P
from PIL import Image, ExifTags

# --- internal imports (unchanged) ---
from .config import DB_FILE, SEMANTIC_METADATA_KEYS
from .path_utils import normalize_path


def _extract_docx(path):
    from .fts_core import extract_virtual_tags

    doc = docx.Document(path)

    clean = (
        lambda x: re.sub(r"\s+", " ", re.sub(r"[\u200b\u00a0\r\n\t]+", " ", x))
        .strip(" .:")
        .strip()
    )

    meta = {}
    for table in doc.tables:
        for row in table.rows:
            if len(row.cells) < 2:
                continue

            key = clean(row.cells[0].text)
            value = clean(row.cells[1].text)

            if key and value and len(value) > 1:
                meta[key.lower()] = value

    paragraphs = [clean(p.text) for p in doc.paragraphs if p.text.strip()]
    full_text = "\n".join(paragraphs)

    extract_virtual_tags(path, text=full_text, meta=meta)
    return full_text


# safe_get helper for .msg and.eml to clean stings


def safe_get(obj, key, fallback=""):
    """Safe getter for dicts, objects, and lists with fallback."""
    try:
        if isinstance(obj, dict):
            value = obj.get(key, fallback)
        else:
            value = getattr(obj, key, fallback)

        if value is None:
            return fallback
        if isinstance(value, (list, tuple)):
            return ", ".join(map(str, value))
        return str(value)
    except Exception:
        return fallback


def _extract_msg(path):
    from .fts_core import extract_virtual_tags

    try:
        msg = extract_msg.Message(path)

        subject = safe_get(msg, "subject", "(no subject)")
        sender = safe_get(msg, "sender", "(unknown sender)")
        to = safe_get(msg, "to", "")
        date = safe_get(msg, "date", "")
        # Prioritize plain > RTF > HTML body
        body = (
            safe_get(msg, "body")
            or safe_get(msg, "bodyRTF")
            or safe_get(msg, "bodyHTML")
        )

        content = f"Subject: {subject}\nFrom: {sender}\nTo: {to}\nDate: {date}\n{body}"
        content = re.sub(r"\s+", " ", content)

        meta = {"From": sender, "To": to, "Subject": subject, "Date": date}
        extract_virtual_tags(path, meta=meta)

        return content.strip()
    except Exception as e:
        print(f"❌ Failed to extract .msg: {e}")
        return ""


def _extract_eml(path):
    from .fts_core import extract_virtual_tags

    try:
        with open(path, "rb") as f:
            raw_email = f.read()

        ep = eml_parser.EmlParser()
        parsed = ep.decode_email_bytes(raw_email)

        header = parsed.get("header", {})
        subject = safe_get(header, "subject", "(no subject)")
        sender = safe_get(header, "from", ["(unknown sender)"])
        to = safe_get(header, "to", [])
        date = safe_get(header, "date", "(no date)")
        body = safe_get(parsed, "body", [""])

        content = f"Subject: {subject}\nFrom: {sender}\nTo: {to}\nDate: {date}\n{body}"
        content = re.sub(r"\s+", " ", content)

        meta = {"From": sender, "To": to, "Subject": subject, "Date": date}
        extract_virtual_tags(path, meta=meta)

        return content.strip()
    except Exception as e:
        print(f"❌ Failed to extract .eml: {e}")
        return ""


def _extract_pptx(path):
    try:
        prs = Presentation(path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text)
    except Exception as e:
        print(f"❌ Failed to extract .pptx: {e}")
        return ""


def _extract_epub(path):
    try:
        book = epub.read_epub(path)
        text = []
        for item in book.get_items():
            if item.get_type() == epub.EpubHtml:
                soup = BeautifulSoup(item.get_body_content(), "html.parser")
                text.append(soup.get_text())
        return "\n".join(text)
    except Exception as e:
        print(f"❌ Failed to extract .epub: {e}")
        return ""


def _extract_odt(path):
    try:
        doc = load(path)
        text = []
        for elem in doc.getElementsByType(P):
            if elem.firstChild:
                text.append(str(elem.firstChild.data))
        return "\n".join(text)
    except Exception as e:
        print(f"❌ Failed to extract .odt: {e}")
        return ""


def _extract_xlsx(path):
    wb = openpyxl.load_workbook(path, data_only=True)
    text = []
    for sheet in wb.worksheets:
        for row in sheet.iter_rows(values_only=True):
            row_text = " ".join([str(cell) for cell in row if cell is not None])
            if row_text.strip():
                text.append(row_text)
    return "\n\n".join(text)


def _extract_pdf(
    path: str,
    ocr_enabled: bool = True,
    force_ocr: bool = False,
    lang: str = "deu+eng",
    max_ocr_pages: int = 3,
    max_pages_for_ocr: int = 10,
    max_size_for_ocr_mb: float = 50.0,
):
    """
    Extract text and metadata from a PDF file.
    Uses PyMuPDF (fitz) for text, with OCR fallback for image-only pages.
    Stores extracted metadata into file_metadata table.
    Smart OCR: skips OCR for large PDFs automatically or when --no-ocr / --ocr is not used.
    """
    text_pages = []
    metadata = {
        "title": None,
        "author": None,
        "subject": None,
        "created": None,
        "last_modified": None,
        "last_modified_by": None,
        "camera": None,
        "image_created": None,
        "dimensions": None,
        "format": "PDF",
        "gps": None,
    }

    try:
        file_size_mb = os.path.getsize(path) / (1024 * 1024)

        with fitz.open(path) as doc:
            num_pages = len(doc)

            # --- Extract metadata from PDF info dictionary ---
            pdf_info = doc.metadata or {}
            metadata.update(
                {
                    "title": pdf_info.get("title"),
                    "author": pdf_info.get("author"),
                    "subject": pdf_info.get("subject"),
                    "created": _normalize_pdf_date(pdf_info.get("creationDate")),
                    "last_modified": _normalize_pdf_date(pdf_info.get("modDate")),
                    "last_modified_by": pdf_info.get("creator")
                    or pdf_info.get("producer"),
                }
            )

            # --- Smart OCR decision ---
            if not ocr_enabled:
                ocr_enabled = False

            elif not force_ocr and (
                num_pages > max_pages_for_ocr or file_size_mb > max_size_for_ocr_mb
            ):
                ocr_enabled = False
                print(
                    f"⚡ Skipping OCR for large PDF "
                    f"({num_pages} pages, {file_size_mb:.1f} MB): {path}"
                )

            elif force_ocr:
                print(f"🔍 OCR enabled (forced): {path}")

            # --- Page text + OCR fallback ---
            for page_num, page in enumerate(doc, start=1):
                page_text = page.get_text("text")

                if page_text.strip():
                    text_pages.append(page_text)
                    continue

                if not ocr_enabled:
                    text_pages.append("")
                    continue

                if not force_ocr and page_num > max_ocr_pages:
                    text_pages.append("")
                    continue

                print(f"📄 OCR page {page_num}/{num_pages}: {path}")

                try:
                    pix = page.get_pixmap(dpi=200)
                    with Image.open(io.BytesIO(pix.tobytes("png"))) as img:
                        ocr_text = pytesseract.image_to_string(img, lang=lang)
                        text_pages.append(ocr_text.strip() if ocr_text else "")
                except Exception as e:
                    print(f"⚠️ OCR failed page {page_num}: {e}")
                    text_pages.append("")

            # --- Fallback to filesystem timestamps ---
            stat = os.stat(path)
            metadata.setdefault(
                "created", datetime.fromtimestamp(stat.st_ctime).isoformat()
            )
            metadata.setdefault(
                "last_modified", datetime.fromtimestamp(stat.st_mtime).isoformat()
            )

            # --- Store metadata in DB ---
            store_metadata(path, metadata)

            # --- Return text and metadata ---
            full_text = "\n\n".join(text_pages).strip()
            return {"text": full_text, "metadata": metadata}

    except Exception as e:
        print(f"⚠️ Error extracting text from {path}: {e}")
        return {"text": "", "metadata": metadata}


def _normalize_pdf_date(date_str):
    """Normalize PDF date strings like 'D:20240512143000Z' to ISO 8601."""
    if not date_str:
        return None
    try:
        date_str = date_str.strip()
        if date_str.startswith("D:"):
            date_str = date_str[2:]
        # Remove timezone or trailing junk
        date_str = date_str.rstrip("Z").split("+")[0].split("-")[0]
        # Parse common PDF formats
        return datetime.strptime(date_str[:14], "%Y%m%d%H%M%S").isoformat()
    except Exception:
        return None


def _extract_html(path):

    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        soup = BeautifulSoup(f.read(), "html.parser")

    # Remove scripts, styles, hidden navs
    for tag in soup(["script", "style", "nav", "footer", "noscript"]):
        tag.decompose()

    # Title
    title = soup.title.string.strip() if soup.title and soup.title.string else ""

    # Extract headings and paragraphs
    elements = []
    for tag in soup.find_all(["h1", "h2", "h3", "h4", "h5", "p"]):
        txt = tag.get_text(" ", strip=True)
        if txt:
            elements.append(txt)

    combined = f"{title}\n" + "\n".join(elements)
    combined = re.sub(r"\s+", " ", combined)  # collapse whitespace
    return combined.strip()


def store_metadata(path, metadata):
    from .db_utils import connect_db

    if not metadata:
        return
    conn = connect_db()
    cursor = conn.cursor()
    columns = [
        "title",
        "author",
        "subject",
        "created",
        "last_modified",
        "last_modified_by",
        "camera",
        "image_created",
        "dimensions",
        "format",
        "gps",
    ]
    values = [metadata.get(col) for col in columns]
    cursor.execute(
        f"""
        INSERT OR REPLACE INTO file_metadata (path, {', '.join(columns)})
        VALUES (?, {', '.join(['?']*len(columns))})
    """,
        [path] + values,
    )
    conn.commit()
    conn.close()


def _gps_to_decimal(coord, ref):
    try:
        d, m, s = coord

        def to_float(x):
            if hasattr(x, "numerator") and hasattr(x, "denominator"):  # IFDRational
                return float(x.numerator) / float(x.denominator)
            elif isinstance(x, tuple):  # (num, den)
                return float(x[0]) / float(x[1])
            elif isinstance(x, (int, float)):  # already float
                return float(x)
            else:
                print(f"⚠️ Unknown EXIF GPS format: {x} ({type(x)})")
                return 0.0

        val = to_float(d) + to_float(m) / 60.0 + to_float(s) / 3600.0
        return -val if ref in ("S", "W") else val

    except Exception as e:
        print(f"⚠️ Failed to parse GPS coord {coord}: {e}")
        return None


def extract_image_metadata(path: str) -> dict:
    """Extract image metadata including EXIF + GPS if present."""
    md = {}
    try:
        with Image.open(path) as img:
            md["dimensions"] = f"{img.width}x{img.height}"
            md["format"] = img.format

            exif = img._getexif()
            if exif:
                exif_data = {ExifTags.TAGS.get(k, k): v for k, v in exif.items()}
                md["camera"] = exif_data.get("Model") or None
                md["image_created"] = exif_data.get("DateTimeOriginal") or None
                md["title"] = exif_data.get("ImageDescription") or None
                md["author"] = exif_data.get("Artist") or None

                gps = exif_data.get("GPSInfo")
                if isinstance(gps, dict):
                    gps_tags = {ExifTags.GPSTAGS.get(k, k): v for k, v in gps.items()}
                    lat = lon = None
                    if all(
                        k in gps_tags
                        for k in (
                            "GPSLatitude",
                            "GPSLatitudeRef",
                            "GPSLongitude",
                            "GPSLongitudeRef",
                        )
                    ):
                        lat = _gps_to_decimal(
                            gps_tags["GPSLatitude"], gps_tags["GPSLatitudeRef"]
                        )
                        lon = _gps_to_decimal(
                            gps_tags["GPSLongitude"], gps_tags["GPSLongitudeRef"]
                        )
                        md["gps"] = f"{lat:.6f},{lon:.6f}"

        stat = os.stat(path)
        md.setdefault("created", datetime.fromtimestamp(stat.st_ctime).isoformat())
        md.setdefault(
            "last_modified", datetime.fromtimestamp(stat.st_mtime).isoformat()
        )
    except Exception as e:
        print(f"⚠️ Failed to extract image metadata: {path} ({e})")
    return md


def update_file_metadata(
    file_path: str,
    metadata: dict,
    conn: sqlite3.Connection | None = None,
) -> str:
    """
    Update the file_metadata table with structured columns and a full JSON column.
    Returns Tier-2-only semantic text (for FTS).

    Internal flags (e.g. content_changed) are persisted but never indexed.
    """
    if not metadata:
        return f"File:{os.path.basename(file_path)}"

    # -------------------------
    # Exclude internal / control flags from FTS semantics
    # -------------------------
    INTERNAL_KEYS = {"content_changed"}

    semantic_parts = [
        f"{k}:{v}"
        for k, v in metadata.items()
        if (
            k in SEMANTIC_METADATA_KEYS
            and k not in INTERNAL_KEYS
            and v not in (None, "", [])
        )
    ]

    semantic_text = (
        " ".join(semantic_parts)
        if semantic_parts
        else f"File:{os.path.basename(file_path)}"
    )

    # -------------------------
    # Structured columns (unchanged)
    # -------------------------
    columns = [
        "title",
        "author",
        "subject",
        "created",
        "last_modified",
        "last_modified_by",
        "camera",
        "image_created",
        "dimensions",
        "format",
        "gps",
    ]
    values = [metadata.get(col) for col in columns]

    close_conn = False
    if conn is None:
        conn = sqlite3.connect(DB_FILE)
        close_conn = True

    try:
        cur = conn.cursor()
        cur.execute(
            f"""
            INSERT OR REPLACE INTO file_metadata
            (path, {', '.join(columns)}, metadata)
            VALUES (?, {', '.join(['?'] * len(columns))}, ?)
            """,
            [file_path, *values, json.dumps(metadata)],
        )

        if close_conn:
            conn.commit()
            conn.close()

    except sqlite3.OperationalError as e:
        print(f"⚠️ Failed to update metadata for {file_path} (SQLite error): {e}")
    except Exception as e:
        print(f"⚠️ Failed to update metadata for {file_path}: {e}")

    return semantic_text

