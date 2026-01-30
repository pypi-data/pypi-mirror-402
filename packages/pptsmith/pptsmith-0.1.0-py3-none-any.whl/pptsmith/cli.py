"""CLI interface for pptsmith."""

import re
from pathlib import Path

import click
from pptx import Presentation


def strip_markdown(text: str) -> str:
    """Remove markdown formatting symbols from text."""
    text = re.sub(r"^\s*[-*+]\s+", "", text, flags=re.MULTILINE)
    text = re.sub(r"\*\*(.+?)\*\*", r"\1", text)
    text = re.sub(r"\*(.+?)\*", r"\1", text)
    text = re.sub(r"__(.+?)__", r"\1", text)
    text = re.sub(r"_(.+?)_", r"\1", text)
    text = re.sub(r"`(.+?)`", r"\1", text)
    text = re.sub(r"~~(.+?)~~", r"\1", text)
    text = re.sub(r"^#{1,6}\s+", "", text, flags=re.MULTILINE)
    return text


def parse_markdown(content: str) -> list[dict]:
    """Parse markdown content into slides."""
    pages = re.split(r"\n---\n", content)
    slides = []

    for page in pages:
        page = page.strip()
        if not page:
            continue

        slide_data = {"title": "", "content": "", "image": None}
        lines = page.split("\n")
        content_lines = []

        for line in lines:
            title_match = re.match(r"^#\s+(.+)$", line)
            if title_match:
                slide_data["title"] = title_match.group(1).strip()
                continue

            img_match = re.match(r"!\[.*?\]\((.+?)\)", line)
            if img_match:
                slide_data["image"] = img_match.group(1)
                continue

            content_lines.append(line)

        content = strip_markdown("\n".join(content_lines).strip())
        # 过滤多个连续空行为一个
        content = re.sub(r'\n{2,}', '\n', content)
        slide_data["content"] = content
        slides.append(slide_data)

    return slides


def fill_slide(slide, slide_data: dict, md_dir: Path, verbose: bool = False, slide_index: int = 0) -> None:
    """Fill a single slide with data, following demo1.py logic exactly."""
    if verbose:
        click.echo(f"\n--- Slide {slide_index + 1} ---")
        click.echo(f"Title: {slide_data['title']}")
        click.echo(f"Content: {slide_data['content'][:100]}..." if len(slide_data['content']) > 100 else f"Content: {slide_data['content']}")
        click.echo(f"Image: {slide_data['image']}")
    # 填充标题占位符 (demo1.py line 9-11)
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 0:
            shape.text = slide_data["title"]

    # 填充文本框内容 - 替换内容但保留样式 (demo1.py line 14-46)
    for shape in slide.shapes:
        if any(kw in shape.name for kw in ["文本框", "Content", "Text"]) and hasattr(shape, "text_frame"):
            tf = shape.text_frame
            if not tf.paragraphs or not tf.paragraphs[0].runs:
                tf.paragraphs[0].text = slide_data["content"]
                continue

            # 保存第一个run的字体样式
            first_run = tf.paragraphs[0].runs[0]
            font_name = first_run.font.name
            font_size = first_run.font.size
            font_bold = first_run.font.bold
            font_italic = first_run.font.italic
            font_color = first_run.font.color.rgb if first_run.font.color.type else None

            # 删除所有段落内容（保留第一个段落的XML结构）
            for para in tf.paragraphs:
                for run in para.runs:
                    run.text = ""
            # 删除多余的段落（只保留第一个）
            p_elements = tf._txBody.findall('.//{http://schemas.openxmlformats.org/drawingml/2006/main}p')
            for p_elem in p_elements[1:]:
                tf._txBody.remove(p_elem)

            # 设置新内容到第一个段落第一个run
            tf.paragraphs[0].runs[0].text = slide_data["content"]
            # 恢复样式
            tf.paragraphs[0].runs[0].font.name = font_name
            if font_size:
                tf.paragraphs[0].runs[0].font.size = font_size
            if font_bold is not None:
                tf.paragraphs[0].runs[0].font.bold = font_bold
            if font_italic is not None:
                tf.paragraphs[0].runs[0].font.italic = font_italic
            if font_color:
                tf.paragraphs[0].runs[0].font.color.rgb = font_color

    # 替换图片，保持宽度适配 (demo1.py line 48-58)
    if slide_data["image"]:
        image_path = md_dir / slide_data["image"]
        if image_path.exists():
            for shape in slide.shapes:
                if any(kw in shape.name for kw in ["图片", "Picture", "Image"]):
                    left = shape.left
                    top = shape.top
                    width = shape.width
                    # 删除原图片
                    sp = shape._element
                    sp.getparent().remove(sp)
                    # 添加新图片，指定宽度（高度自动按比例计算）
                    slide.shapes.add_picture(str(image_path), left, top, width=width)
                    break
    else:
        # 没有图片时删除模板中的图片占位
        for shape in slide.shapes:
            if any(kw in shape.name for kw in ["图片", "Picture", "Image"]):
                sp = shape._element
                sp.getparent().remove(sp)
                break


def duplicate_slide(prs, slide):
    """Duplicate a slide by copying its XML structure."""
    import copy
    from pptx.util import Inches
    
    # 获取 slide 的 XML 并深拷贝
    template_xml = slide._element
    new_slide_xml = copy.deepcopy(template_xml)
    
    # 添加新 slide 到 presentation
    slide_layout = slide.slide_layout
    new_slide = prs.slides.add_slide(slide_layout)
    
    # 清空新 slide 的所有元素
    spTree = new_slide.shapes._spTree
    for child in list(spTree):
        if child.tag.endswith('}sp') or child.tag.endswith('}pic'):
            spTree.remove(child)
    
    # 从模板复制所有元素
    template_spTree = template_xml.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}spTree')
    for child in template_spTree:
        if child.tag.endswith('}sp') or child.tag.endswith('}pic'):
            spTree.append(copy.deepcopy(child))
    
    return new_slide


def create_presentation(
    slides_data: list[dict], template_path: Path, output_path: Path, md_dir: Path, verbose: bool = False
) -> None:
    """Create PowerPoint presentation from parsed slides data."""
    import copy
    
    if not slides_data:
        return
    
    prs = Presentation(str(template_path))
    template_slide = prs.slides[0]
    
    # 先复制出所需数量的 slide（从第2页开始）
    slides_to_fill = [template_slide]
    for _ in range(len(slides_data) - 1):
        new_slide = duplicate_slide(prs, template_slide)
        slides_to_fill.append(new_slide)
    
    # 逐个填充
    for i, (slide, slide_data) in enumerate(zip(slides_to_fill, slides_data)):
        fill_slide(slide, slide_data, md_dir, verbose, i)

    prs.save(str(output_path))


@click.command()
@click.argument("markdown_file", type=click.Path(exists=True, path_type=Path))
@click.option(
    "-t",
    "--template",
    type=click.Path(exists=True, path_type=Path),
    default="template.pptx",
    help="PowerPoint template file (default: template.pptx)",
)
@click.option(
    "-o",
    "--output",
    type=click.Path(path_type=Path),
    default=None,
    help="Output file path (default: same name as input with .pptx extension)",
)
@click.option(
    "-v",
    "--verbose",
    is_flag=True,
    default=False,
    help="Print rendering content for each slide",
)
def main(markdown_file: Path, template: Path, output: Path, verbose: bool) -> None:
    """Convert Markdown file to PowerPoint presentation.

    MARKDOWN_FILE: Path to the markdown file to convert.
    """
    if output is None:
        output = markdown_file.with_suffix(".pptx")

    content = markdown_file.read_text(encoding="utf-8")
    slides_data = parse_markdown(content)

    md_dir = markdown_file.parent

    create_presentation(slides_data, template, output, md_dir, verbose)

    click.echo(f"Created presentation: {output}")
    click.echo(f"Total slides: {len(slides_data)}")


if __name__ == "__main__":
    main()
